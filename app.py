"""app.py — Impact Project Room (FULL SCRIPT)
================================================
Entrepreneurs upload a confidential PDF + metadata. GPT‑4o extracts nine
fields into JSON. Admins browse, filter, update status and export.
"""

import os
from pathlib import Path
from dotenv import load_dotenv

import streamlit as st

 
 # ---- Pandas Styler helpers (table colouring) -------------------
def style_status(val):
    color = STATUS_COLORS.get(val, "#ffffff")
    text = "#fff" if int(color.lstrip("#"), 16) < 0x888888 else "#000"
    return f"background-color: {color}; color:{text}"

def style_region(val):
    color = REGION_COLORS.get(val, "#ffffff")
    return f"background-color: {color};"

def style_sector(val):
    color = SECTOR_COLORS.get(val, "#ffffff")
    return f"background-color: {color};"

import pydeck as pdk
# Configure page layout
st.set_page_config(layout="wide", page_title="Impact Project Room")
# ── Brand Color Theme ──────────────────────────────────────────
# Custom brand colors for buttons, headers, sidebar, and charts
st.markdown(
    """
    <style>
    /* Primary color (Orange) on buttons, download buttons, tabs, and expanders */
    .stButton button,
    .stDownloadButton button,
    .stTabs button,
    .stExpanderHeader {
        background-color: #D97A45 !important;
        border-color: #D97A45 !important;
        color: #ffffff !important;
    }
    /* Secondary color (Dark Blue) on headers and metric titles */
    h1, h2, h3, .stMetric {
        color: #243D66 !important;
    }
    /* Tertiary color (Light Blue) for sidebar background */
    .css-1lcbmhc, .css-1avcm0n {
        background-color: #516f98 !important;
    }
    </style>
    """,
    unsafe_allow_html=True
)
import openai

# ── Imports & basic setup ──────────────────────────────────────────
import os
import io
import json
from datetime import datetime
from typing import Dict, Any
import difflib
import random

import fitz  # PyMuPDF
import openai
import pandas as pd
import streamlit as st
from docx import Document
from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib as mpl
mpl.rcParams['axes.prop_cycle'] = mpl.cycler(color=['#D97A45', '#243D66', '#516f98'])
try:
    import plotly.express as px
    HAS_PLOTLY = True
except ImportError:
    HAS_PLOTLY = False
    px = None
# For dynamic country centroids
try:
    from countryinfo import CountryInfo
    HAS_COUNTRYINFO = True
except ImportError:
    HAS_COUNTRYINFO = False
    # Fallback stub so mapping code can run without crashing
    class CountryInfo:
        def __init__(self, name):
            pass
        @property
        def latlng(self):
            return (None, None)
from PIL import Image

from pathlib import Path
from dotenv import load_dotenv

# Load .env locally
load_dotenv(Path(__file__).parent / ".env")


# Load local .env for development
dotenv_path = Path(__file__).parent / ".env"
load_dotenv(dotenv_path)

# Load .env for local development
from dotenv import load_dotenv
load_dotenv()  # load .env for local development

# --- utilities -----------------------------------------------------
# ----- simple debug toggle ---------------------------------------
# Always show GPT raw output in logs and UI.
DEBUG_GPT = True  # always show GPT raw output

def safe_int(x, default):
    """
    Convert x to int, or return `default` if x is None, empty, or not an int-ish string.
    """
    try:
        return int(x)
    except (TypeError, ValueError):
        return default
    
# ── Environment & OpenAI setup ────────────────────────────────────
from pathlib import Path

# Load app logo
logo_img = Image.open(Path(__file__).with_name("logo.png"))

# SharePoint integration
from office365.sharepoint.client_context import ClientContext
from office365.runtime.auth.client_credential import ClientCredential


# Load OpenAI API key: prefer Streamlit secrets, then environment (.env), then .env local file
api_key = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
if not api_key:
    st.error(
        "🔒 OPENAI_API_KEY not found.\n"
        "- For local dev: create a .env file with OPENAI_API_KEY or set the env var.\n"
        "- For Streamlit Cloud: add it under App Settings → Secrets."
    )
    st.stop()
openai.api_key = api_key

# Show logo at top
# Show small logo at top-right
col1, col2 = st.columns([9, 1])
with col2:
    st.image(logo_img, width=60)

UPLOAD_FOLDER = "submissions"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

AI_FIELDS = [
    "Specific Sector(s)",
    "Region of operation",                 # Global, Western Economies, Africa, Asia, SEA, Latam
    "Main country of current operations",
    "Business Model",
    "Maturity stage",
    "Core team",                               # leadership team
    "Key risks",
    "Last 12 months revenues (USD)",
    "Breakeven year",
    "Market size or SOM (USD)",
    "Expected IRR (%)",
    "Financing need or round size (USD)",
    "Instrument",
    "Use of proceeds",
    "Impact Area",
    "Main SDGs targeted (3 max)",
    "Problem",
    "Solution",
    "Barrier(s) to entry",
]
# Metadata fields manually provided by entrepreneur/admin
METADATA_FIELDS = [
    "Project registered name",
    "Date of incorporation",
    "Primary sector / theme",
    "Headquarters country",
    "Contact e-mail",
]

# ── Field type helpers ────────────────────────────────────────────
NUMERIC_USD_FIELDS = [
    "Last 12 months revenues (USD)",
    "Market size or SOM (USD)",
    "Financing need or round size (USD)",
]
NUMERIC_PERCENT_FIELDS = [
    "Expected IRR (%)",
]
YEAR_FIELDS = ["Breakeven year"]

REVIEW_STAGES = [
    "Identified",
    "Intro call",
    "NDA and Deck",
    "Financials",
    "4-pager",
    "IC1",
    "IC2",
    "Local DD",
    "Raised",
    "Operating",
    "Exited",
    "Bankrupt",
]

# Color mapping for Due Diligence stages (light to dark blue)
STATUS_COLORS = {
    "Identified": "#deebf7",
    "Intro call": "#c6dbef",
    "NDA and Deck": "#9ecae1",
    "Financials": "#6baed6",
    "4-pager": "#4292c6",
    "IC1": "#2171b5",
    "IC2": "#08519c",
    "Local DD": "#084594",
    "Raised": "#2171b5",
    "Operating": "#08519c",
    "Exited": "#041f4a",
    "Bankrupt": "#020f1f",
}

# Color mapping for Regions
REGION_COLORS = {
    "Global": "#808080",
    "Western Economies": "#00aaff",
    "Africa": "#008000",
    "Asia": "#FFA500",
    "SEA": "#A52A2A",
    "Latam": "#FF0000",
}

# Color mapping for Sectors
SECTOR_COLORS = {
    "Agriculture": "#2ca02c",
    "Air": "#17becf",
    "Biodiversity & ecosystems": "#9467bd",
    "Climate": "#7f7f7f",
    "Diversity & inclusion": "#e377c2",
    "Education": "#8c564b",
    "Employment / Livelihoods creation": "#bcbd22",
    "Energy": "#d62728",
    "Financial services": "#1f77b4",
    "Health": "#ff7f0e",
    "Infrastructure": "#7f7f7f",
    "Land": "#8c564b",
    "Oceans & coastal zones": "#17becf",
    "Sustainable cities": "#2ca02c",
    "Sustainable consumption & production": "#bcbd22",
    "Sustainable tourism": "#e377c2",
    "Water Treatment": "#1f77b4",
    "Other": "#c7c7c7",
}

# Color mapping for Financing Instrument types
INSTRUMENT_COLORS = {
    "convertible note": "#6a3d9a",
    "equity": "#e31a1c",
    "debt": "#1f78b4",
    "other": "#b15928",
}

# Options for the entrepreneur submission form
SECTOR_OPTIONS = [
    "Agriculture",
    "Air",
    "Biodiversity & ecosystems",
    "Climate",
    "Diversity & inclusion",
    "Education",
    "Employment / Livelihoods creation",
    "Energy",
    "Financial services",
    "Health",
    "Infrastructure",
    "Land",
    "Oceans & coastal zones",
    "Sustainable cities",
    "Sustainable consumption & production",
    "Sustainable tourism",
    "Water Treatment",
    "Other",
]

# ISO‑style country display list 
COUNTRY_OPTIONS = [
    "Afghanistan","Albania","Algeria","American Samoa","Andorra","Angola",
    "Anguilla","Antarctica","Antigua And Barbuda","Argentina","Armenia",
    "Aruba","Australia","Austria","Azerbaijan","Bahamas The","Bahrain",
    "Bangladesh","Barbados","Belarus","Belgium","Belize","Benin","Bermuda",
    "Bhutan","Bolivia","Bosnia and Herzegovina","Botswana","Bouvet Island",
    "Brazil","British Indian Ocean Territory","Brunei","Bulgaria",
    "Burkina Faso","Burundi","Cambodia","Cameroon","Canada","Cape Verde",
    "Cayman Islands","Central African Republic","Chad","Chile","China",
    "Christmas Island","Cocos (Keeling) Islands","Colombia","Comoros",
    "Republic Of The Congo","Democratic Republic Of The Congo","Cook Islands",
    "Costa Rica","Cote D'Ivoire (Ivory Coast)","Croatia (Hrvatska)","Cuba",
    "Cyprus","Czech Republic","Denmark","Djibouti","Dominica",
    "Dominican Republic","East Timor","Ecuador","Egypt","El Salvador",
    "Equatorial Guinea","Eritrea","Estonia","Ethiopia",
    "External Territories of Australia","Falkland Islands","Faroe Islands",
    "Fiji Islands","Finland","France","French Guiana","French Polynesia",
    "French Southern Territories","Gabon","Gambia The","Georgia","Germany",
    "Ghana","Gibraltar","Greece","Greenland","Grenada","Guadeloupe","Guam",
    "Guatemala","Guernsey and Alderney","Guinea","Guinea-Bissau","Guyana",
    "Haiti","Heard and McDonald Islands","Honduras","Hong Kong S.A.R.",
    "Hungary","Iceland","India","Indonesia","Iran","Iraq","Ireland","Israel",
    "Italy","Jamaica","Japan","Jersey","Jordan","Kazakhstan","Kenya",
    "Kiribati","Korea North","Korea South","Kuwait","Kyrgyzstan","Laos",
    "Latvia","Lebanon","Lesotho","Liberia","Libya","Liechtenstein",
    "Lithuania","Luxembourg","Macau S.A.R.","Macedonia","Madagascar",
    "Malawi","Malaysia","Maldives","Mali","Malta","Marshall Islands",
    "Martinique","Mauritania","Mauritius","Mayotte","Mexico","Micronesia",
    "Moldova","Monaco","Mongolia","Montenegro","Montserrat","Morocco",
    "Mozambique","Myanmar (Burma)","Namibia","Nauru","Nepal","Netherlands",
    "Netherlands Antilles","New Caledonia","New Zealand","Nicaragua","Niger",
    "Nigeria","Niue","Norfolk Island","Northern Mariana Islands","Norway",
    "Oman","Pakistan","Palau","Palestinian Territories","Panama",
    "Papua New Guinea","Paraguay","Peru","Philippines","Pitcairn Islands",
    "Poland","Portugal","Puerto Rico","Qatar","Reunion","Romania","Russia",
    "Rwanda","Saint Helena","Saint Kitts and Nevis","Saint Lucia",
    "Saint Pierre and Miquelon","Saint Vincent and the Grenadines","Samoa",
    "San Marino","Sao Tome and Principe","Saudi Arabia","Senegal","Serbia",
    "Seychelles","Sierra Leone","Singapore","Slovakia","Slovenia",
    "Solomon Islands","Somalia","South Africa",
    "South Georgia and the South Sandwich Islands","Spain","Sri Lanka",
    "Sudan","Suriname","Svalbard and Jan Mayen","Swaziland","Sweden",
    "Switzerland","Syria","Taiwan","Tajikistan","Tanzania","Thailand",
    "Timor-Leste (East Timor)","Togo","Tokelau","Tonga","Trinidad and Tobago",
    "Tunisia","Turkey","Turkmenistan","Turks and Caicos Islands","Tuvalu",
    "Uganda","Ukraine","United Arab Emirates","United Kingdom",
    "United States","United States Minor Outlying Islands","Uruguay",
    "Uzbekistan","Vanuatu","Vatican City","Venezuela","Vietnam",
    "Virgin Islands, British","Virgin Islands, U.S.","Wallis and Futuna",
    "Western Sahara","Yemen","Zambia","Zimbabwe",
]


INSTRUMENT_OPTIONS = ["Convertible note", "Equity", "Debt", "Other"]

# Portfolio holder tags (added by NCGE / NCGD interfaces)
PORTFOLIO_OPTIONS = ["NCGE", "NCGD"]

# ── Standard options for review-stage input ───────────────────────
SDG_OPTIONS = [
    "No poverty (SDG 1)",
    "Zero hunger (SDG 2)",
    "Good health and well-being (SDG 3)",
    "Quality education (SDG 4)",
    "Gender equality (SDG 5)",
    "Clean water and sanitation (SDG 6)",
    "Affordable and clean energy (SDG 7)",
    "Decent work and economic growth (SDG 8)",
    "Industry, innovation and infrastructure (SDG 9)",
    "Reduced inequalities (SDG 10)",
    "Sustainable cities and communities (SDG 11)",
    "Responsible consumption and production (SDG 12)",
    "Climate action (SDG 13)",
    "Life below water (SDG 14)",
    "Life on land (SDG 15)",
    "Peace, justice, and strong institutions (SDG 16)",
    "Partnerships for the goals (SDG 17)",
]

MATURITY_STAGES = [
    "Ideation",
    "Validation",
    "Pilot",
    "Growth",
    "Scale",
    "Mature",
]

# Granular list used for multi-select “Specific Sector(s)”
SPECIFIC_SECTOR_OPTIONS = [
    "AgTech", "CleanTech", "EdTech", "HealthTech", "FinTech", "ClimateTech",
    "Water & Sanitation", "Circular Economy", "Energy Access", "Sustainable Forestry",
    "Sustainable Agriculture", "Waste Management", "Wildlife Conservation",
    "Blue Economy", "Mobility & Transport", "Carbon Markets", "Green Buildings",
    "E-commerce for Impact"
]

# Fuzzy-matching helper for SDGs
def _match_sdgs(raw_list):
    """
    Fuzzy-match a list of raw SDG strings to the canonical SDG_OPTIONS.
    Returns up to 3 unique best matches with cutoff=0.6.
    """
    matched = []
    for s in raw_list:
        candidates = difflib.get_close_matches(s, SDG_OPTIONS, n=1, cutoff=0.6)
        if candidates and candidates[0] not in matched:
            matched.append(candidates[0])
        if len(matched) >= 3:
            break
    return matched

# ── Utility helpers ───────────────────────────────────────────────

# Tiny Unicode sparklines (▁ ▂ ▃ ▄ ▅ ▆ ▇ █)
def _sparkline(arr):
    """
    Return a compact sparkline string from a list of numbers.
    Example: [1,2,3] -> '▁▄█'
    """
    ticks = "▁▂▃▄▅▆▇█"
    if not arr:
        return ""
    mn, mx = min(arr), max(arr)
    span = mx - mn if mx != mn else 1e-9
    return "".join(ticks[int((v - mn) / span * (len(ticks) - 1))] for v in arr)

def extract_text_from_pdf(path: str) -> str:
    """Return concatenated text from a PDF, trying multiple extraction methods."""
    doc = fitz.open(path)
    chunks = []
    for page in doc:
        txt = page.get_text("text")  # fastest
        if not txt.strip():
            # Some PDFs render better with the dict API
            try:
                dict_blocks = page.get_text("dict")["blocks"]
                for block in dict_blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                if span.get("text"):
                                    chunks.append(span["text"])
            except Exception:
                pass
        else:
            chunks.append(txt)
    return "\n".join(chunks)


# Extract text from various file types (PDF, DOCX, PPTX, XLSX)
def extract_text_from_file(path: str) -> str:
    """
    Extract text from a PDF, DOCX, PPTX, or XLSX file based on its extension.
    """
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".docx":
        doc = Document(path)
        return "\n".join(para.text for para in doc.paragraphs)
    elif ext == ".pptx":
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    elif ext in [".xls", ".xlsx"]:
        try:
            df = pd.read_excel(path)
            return df.to_csv(index=False)
        except Exception:
            return ""
    else:
        return ""



def _parse_json_from_string(payload: str) -> Dict[str, Any]:
    payload = payload.strip()
    try:
        return json.loads(payload)
    except Exception:
        if "{" in payload and "}" in payload:
            snippet = payload[payload.find("{") : payload.rfind("}") + 1]
            try:
                return json.loads(snippet)
            except Exception:
                pass
    return {}

# ── Helper: load layered summary (entrepreneur/admin) ────────────
def load_layered_summary(folder: str):
    """
    Returns (summary_dict, audit_dict) with an added flag 'entrepreneur_pending' in audit_dict.
    Logic:
    1. Use admin summary if present; if entrepreneur edits are newer, keep admin but set pending flag.
    2. Else use entrepreneur summary if present.
    3. Else fallback to GPT summary.
    """
    import os
    from datetime import datetime

    # Paths
    ent_json = os.path.join(folder, "summary_entrepreneur.json")
    ent_txt  = os.path.join(folder, "summary_gpt.txt")
    adm_json = os.path.join(folder, "summary_admin.json")

    # Load entrepreneur version (JSON preferred, then GPT TXT)
    ent, ent_audit = {}, {}
    ent_mtime = None
    if os.path.exists(ent_json):
        data = json.load(open(ent_json, encoding="utf-8"))
        if isinstance(data, dict) and "fields" in data:
            ent = data["fields"]
            ent_audit = data.get("audit", {})
        else:
            ent = data
        ent_mtime = os.path.getmtime(ent_json)
    elif os.path.exists(ent_txt):
        raw = open(ent_txt, encoding="utf-8").read()
        ent = _parse_json_from_string(raw)
        ent_mtime = os.path.getmtime(ent_txt)

    # Load admin version if exists
    adm, adm_audit = {}, {}
    adm_mtime = None
    if os.path.exists(adm_json):
        try:
            data = json.load(open(adm_json, encoding="utf-8"))
            adm = data.get("fields", {})
            adm_audit = data.get("audit", {})
            adm_mtime = os.path.getmtime(adm_json)
        except json.JSONDecodeError:
            adm, adm_audit = {}, {}

    # Determine current summary and entrepreneur_pending flag
    entrepreneur_pending = False
    if adm_mtime:
        # Admin exists: always use admin content
        summary = adm
        audit = adm_audit
        # But if entrepreneur version exists and is newer, flag pending edits
        if ent_mtime and ent_mtime > adm_mtime:
            entrepreneur_pending = True
    elif ent_mtime:
        # No admin: use entrepreneur version
        summary = ent
        audit = ent_audit
    else:
        # Fallback GPT-only (ent must have loaded from GPT)
        summary = ent
        audit = ent_audit

    # Attach pending flag
    audit["entrepreneur_pending"] = entrepreneur_pending
    return summary, audit


def summarize_project_with_gpt(full_text: str) -> Dict[str, Any]:
    """
    Extract structured impact project fields from text using a direct JSON prompt.
    Falls back missing keys to "Unknown".
    """
    # Truncate to fit context window
    text = full_text[:25000]  # increased context window

    # Build a bullet‑point style system prompt for better adherence
    schema_keys = ", ".join([f'"{k}"' for k in AI_FIELDS])
    system_prompt = (
        "You are an expert impact‑investment analyst.\n"
        "Return **exactly** the following JSON schema, keys in this order and never omit any key "
        f"({schema_keys}). If a value is missing, use an empty string \"\" or empty list [].\n"
        "Guidelines:\n"
        "• Use the **closest match** from each allowed list (case‑insensitive). For example, \"Convertible Note\" should map to \"Convertible note\".\n"
        "• The allowed lists are:\n"
        f"  – Specific Sector(s): {', '.join(SECTOR_OPTIONS)}\n"
        f"  – Region of operation: {', '.join(REGION_COLORS.keys())}\n"
        f"  – Main country of current operations: {', '.join(COUNTRY_OPTIONS)}\n"
        f"  – Maturity stage: {', '.join(MATURITY_STAGES)}\n"
        f"  – Instrument: {', '.join(INSTRUMENT_OPTIONS)}\n"
        f"  – Main SDGs targeted: {', '.join(SDG_OPTIONS)}\n"
        "• If it does not mention relevant information for a field, leave it empty – unless you believe you can confidently guess it.\n"
        "Output **only** the JSON (no commentary)."
    )
    user_prompt = f"Pitch Content:\n{text}"

    try:
        resp = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.0,
            max_tokens=2000,
        )
        raw_output = resp.choices[0].message.content.strip()
        # Store raw output so it can be shown elsewhere in the UI
        st.session_state["last_raw_output"] = raw_output
        st.session_state["last_raw_output"] = raw_output
        if DEBUG_GPT:
            print("\n===== GPT RAW OUTPUT =====")
            print(raw_output)
            print("===== END =====\n")
            # Display the raw JSON in the Streamlit UI for easy copy‑paste
            with st.expander("🔍 GPT raw output", expanded=False):
                st.code(raw_output, language="json")
        # Parse JSON payload from the response
        summary = _parse_json_from_string(raw_output)
    except Exception as e:
        st.error(f"OpenAI API error: {e}")
        summary = {}

    # Ensure every expected field exists (skip when debugging to see true payload)
    if not DEBUG_GPT:
        for key in AI_FIELDS:
            summary.setdefault(key, "Unknown")

    return summary


def render_summary_grid(summary: Dict[str, Any]):
    # Combine metadata and AI fields for display
    fields = METADATA_FIELDS + [f for f in AI_FIELDS if f not in METADATA_FIELDS]
    cols = st.columns(3)
    for idx, field in enumerate(fields):
        val = summary.get(field, "–")
        if isinstance(val, list):
            val = "; ".join(
                f"{m.get('Name','')} ({m.get('Role','')})" if isinstance(m, dict) else str(m)
                for m in val
            )
        cols[idx % 3].markdown(f"**{field}**  \n{val}")

# ── Helper: save the submission to disk ───────────────────────────
def _save_submission(meta: dict, files, summary: dict):
    """
    Creates a timestamped folder in /submissions and writes:
    • each uploaded file
    • info.txt with human‑readable metadata
    • summary_entrepreneur.json containing the AI JSON
    """
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_proj = meta["Project"].replace(" ", "_")
    fld = os.path.join(UPLOAD_FOLDER, f"{safe_proj}_{ts}")
    os.makedirs(fld, exist_ok=True)

    # Save uploaded files
    for upl in files:
        fpath = os.path.join(fld, upl.name)
        with open(fpath, "wb") as f:
            f.write(upl.read())

    # Write metadata
    meta_text = "\n".join(f"{k}: {v}" for k, v in meta.items()) + "\nNDA: Accepted\n"
    with open(os.path.join(fld, "info.txt"), "w", encoding="utf-8") as mf:
        mf.write(meta_text)

    # Save entrepreneur AI summary with audit + empty history
    ent_path = os.path.join(fld, "summary_entrepreneur.json")
    audit = {"entrepreneur_last_updated": datetime.now().isoformat()}
    with open(ent_path, "w", encoding="utf-8") as ef:
        json.dump({"fields": summary, "audit": audit, "history": []},
              ef, ensure_ascii=False, indent=2)
    # Mirror submission folder to SharePoint
    try:
        _upload_to_sharepoint(fld, os.path.basename(fld))
    except Exception as e:
        print("SharePoint upload error:", e)

    # Generate and save edit PIN for the entrepreneur
    pin = f"{random.randint(0, 9999):04d}"
    cred_path = os.path.join(fld, "credentials.json")
    with open(cred_path, "w", encoding="utf-8") as cf:
        json.dump({"pin": pin}, cf)
    return fld, pin


# ── Helper: update an existing submission ────────────────────────
def _update_submission(folder: str, meta: dict, summary: dict):
    base = os.path.join(UPLOAD_FOLDER, folder)

    # -- info.txt (metadata) -------------------------------------------------
    meta_text = "\n".join(f"{k}: {v}" for k, v in meta.items()) + "\nNDA: Accepted\n"
    with open(os.path.join(base, "info.txt"), "w", encoding="utf-8") as mf:
        mf.write(meta_text)

    # -- versioned entrepreneur summary --------------------------------------
    ent_path = os.path.join(base, "summary_entrepreneur.json")
    new_audit = {"entrepreneur_last_updated": datetime.now().isoformat()}

    history = []
    if os.path.exists(ent_path):
        try:
            prev = json.load(open(ent_path, encoding="utf-8"))
            history = prev.get("history", [])
            history.insert(0, {
                "fields": prev.get("fields", prev),
                "audit":  prev.get("audit", {})
            })
        except Exception:
            pass  # corrupt file → start fresh

    with open(ent_path, "w", encoding="utf-8") as ef:
        json.dump({"fields": summary, "audit": new_audit, "history": history},
                ef, ensure_ascii=False, indent=2)

# ── Helper: save admin summary (for table edits) ────────────────
def _save_admin_summary(folder: str, fields: Dict[str, Any]):
    """
    Overwrite *summary_admin.json* with the latest admin‑curated fields **and**
    maintain a full version history.

    The file structure becomes:
    {
        "fields": {...},          # latest admin version used by the app
        "audit":  {"admin_last_updated": "<ISO‑timestamp>"},
        "history": [              # older admin versions, newest first
            {
                "fields": {...},
                "audit":  {...}
            },
            ...
        ]
    }
    """
    base = os.path.join(UPLOAD_FOLDER, folder)
    adm_path = os.path.join(base, "summary_admin.json")
    new_audit = {"admin_last_updated": datetime.now().isoformat()}

    # If a previous admin summary exists, load it and move it to history
    history: list = []
    if os.path.exists(adm_path):
        try:
            with open(adm_path, encoding="utf-8") as f:
                prev_data = json.load(f)
            # Push the current content to history (ignore if corrupt)
            history = prev_data.get("history", [])
            history.insert(0, {              # newest first
                "fields": prev_data.get("fields", {}),
                "audit":  prev_data.get("audit", {})
            })
        except Exception:
            pass  # corrupt file – start fresh

    # Write the new file
    with open(adm_path, "w", encoding="utf-8") as af:
        json.dump({"fields": fields, "audit": new_audit, "history": history},
                  af, ensure_ascii=False, indent=2)


# --- Portfolio helper: add/eject portfolio tag in summary_admin.json ---
def _set_portfolio_membership(folder: str, role_tag: str, action: str):
    """
    Manage the role_tag in Portfolio / Rejected lists.

    * action = "add"        → ensure tag in Portfolio, remove from Rejected
    * action = "eject"      → ensure tag in Rejected, remove from Portfolio
    * action = "reconsider" → remove tag from Rejected (if present)
    """
    base = os.path.join(UPLOAD_FOLDER, folder)
    admin_path = os.path.join(base, "summary_admin.json")

    # ---------- Load existing fields ----------
    if os.path.exists(admin_path):
        try:
            adm = json.load(open(admin_path, encoding="utf-8"))
            fields = adm.get("fields", {})
        except Exception:
            fields = {}
    else:
        # Fallback: start from entrepreneur version
        ent_path = os.path.join(base, "summary_entrepreneur.json")
        if os.path.exists(ent_path):
            try:
                ent_raw = json.load(open(ent_path, encoding="utf-8"))
                fields = ent_raw.get("fields", ent_raw) if isinstance(ent_raw, dict) else ent_raw
            except Exception:
                fields = {}
        else:
            fields = {}

    # Normalise lists
    port = fields.get("Portfolio", [])
    rej  = fields.get("Rejected", [])
    if isinstance(port, str):
        port = [p.strip() for p in port.split(";") if p.strip()]
    if isinstance(rej, str):
        rej = [p.strip() for p in rej.split(";") if p.strip()]

    # ---------- Apply action ----------
    if action == "add":
        if role_tag not in port:
            port.append(role_tag)
        if role_tag in rej:
            rej.remove(role_tag)
    elif action == "eject":
        if role_tag not in rej:
            rej.append(role_tag)
        if role_tag in port:
            port.remove(role_tag)
    elif action == "reconsider":
        if role_tag in rej:
            rej.remove(role_tag)

    # ---------- Save back ----------
    fields["Portfolio"] = port
    fields["Rejected"]  = rej
    _save_admin_summary(folder, fields)

def _upload_to_sharepoint(local_folder: str, project_folder: str):
    """
    Upload all files from local_folder to a SharePoint document library,
    creating a subfolder named project_folder.
    """
    # Pull SharePoint config from secrets
    site_url      = st.secrets["SP_SITE_URL"]
    client_id     = st.secrets["SP_CLIENT_ID"]
    client_secret = st.secrets["SP_CLIENT_SECRET"]
    library       = st.secrets["SP_DOC_LIBRARY"]  # e.g. "Shared Documents/ImpactSubmissions"

    # Authenticate
    creds = ClientCredential(client_id, client_secret)
    ctx   = ClientContext(site_url).with_credentials(creds)

    # Create (or get) the project folder
    root_folder      = ctx.web.get_folder_by_server_relative_url(library)
    project_sp_folder = root_folder.add_folder(project_folder).execute_query()

    # Upload each file in the local folder
    for fname in os.listdir(local_folder):
        path = os.path.join(local_folder, fname)
        with open(path, "rb") as f:
            content = f.read()
        project_sp_folder.upload_file(fname, content).execute_query()

# ── Helper: rerun compatible with all Streamlit versions ──────────
def _rerun():
    if hasattr(st, "rerun"):            # Streamlit ≥ 1.27
        st.rerun()
    else:                               # older versions
        st.experimental_rerun()

# ── Helper: send e‑mail alert when a new submission arrives ──────
import smtplib, ssl
from email.message import EmailMessage

def email_admin(subject: str, body: str):
    """Send a plain‑text e‑mail to the address set in secrets or fallback literal."""
    host = st.secrets.get("SMTP_HOST")
    port = int(st.secrets.get("SMTP_PORT", 0))
    user = st.secrets.get("SMTP_USER")
    pwd  = st.secrets.get("SMTP_PASS")
    to   = st.secrets.get("ADMIN_EMAIL", "NCM@kickimpact.com")

    if not all([host, port, user, pwd]):
        print("Email credentials missing: alert not sent.")
        return

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"]    = user
    msg["To"]      = to
    msg.set_content(body)

    ctx = ssl.create_default_context()
    try:
        if port == 465:   # SSL
            with smtplib.SMTP_SSL(host, port, context=ctx) as srv:
                srv.login(user, pwd)
                srv.send_message(msg)
        else:             # STARTTLS
            with smtplib.SMTP(host, port) as srv:
                srv.starttls(context=ctx)
                srv.login(user, pwd)
                srv.send_message(msg)
    except Exception as e:
        print("E‑mail error:", e)

# ── Routing ───────────────────────────────────────────────────────
is_admin = str(st.query_params.get("adminNCM", "")).lower() == "true"
 # Detect Portfolio‑holder roles (Gross Equity / Gross Debt)
is_portfolio_ncge = str(st.query_params.get("adminNCGE", "")).lower() == "true"
is_portfolio_ncgd = str(st.query_params.get("adminNCGD", "")).lower() == "true"

#
# ═════════════════════════════════════════════════════════════════
# ADMIN DASHBOARD
# ═════════════════════════════════════════════════════════════════
if is_admin:
    st.title("🛠️ Admin Dashboard")
    # Sidebar filter panel (scrollable)
    st.sidebar.header("🔍 Filters")
    filter_logic = st.sidebar.radio(
        "Filter logic",
        ["AND", "OR"],
        index=0,
        help="Use AND to require all selected filters, OR to include any matching filter."
    )

    # Categorical Filters
    st.sidebar.markdown("#### Categorical Filters")
    # Primary sector / theme
    primary_sector_filter = st.sidebar.multiselect(
        "Primary sector / theme", SECTOR_OPTIONS
    )
    # Headquarters country
    hq_country_filter = st.sidebar.multiselect(
        "Headquarters country", COUNTRY_OPTIONS
    )
    # Specific Sector(s)
    specific_sector_filter = st.sidebar.multiselect(
        "Specific Sector(s)", SPECIFIC_SECTOR_OPTIONS
    )
    # Region of operation
    region_filter = st.sidebar.multiselect(
        "Region of operation", list(REGION_COLORS.keys())
    )
    # Main country of current operations
    main_country_filter = st.sidebar.multiselect(
        "Main country of current operations", COUNTRY_OPTIONS
    )
    # Instrument
    instrument_filter = st.sidebar.multiselect(
        "Instrument", INSTRUMENT_OPTIONS
    )
    # Maturity stage
    maturity_filter = st.sidebar.multiselect(
        "Maturity stage", MATURITY_STAGES
    )
    # Main SDGs targeted (3 max)
    sdg_filter = st.sidebar.multiselect(
        "Main SDGs targeted (3 max)", SDG_OPTIONS
    )
    # Due Diligence stage
    stage_filter = st.sidebar.multiselect(
        "Due Diligence stage", REVIEW_STAGES
    )

    # Numerical filters (bar-style sliders)
    st.sidebar.markdown("#### Numerical Filters")
    irr_min, irr_max = st.sidebar.slider(
        "IRR Range (%)", min_value=0.0, max_value=100.0, value=(0.0, 100.0), step=0.1
    )
    revenue_min, revenue_max = st.sidebar.slider(
        "Revenue Range (USD)", min_value=0, max_value=10_000_000, value=(0, 10_000_000), step=1_000
    )
    som_min, som_max = st.sidebar.slider(
        "SOM Range (USD)", min_value=0, max_value=100_000_000, value=(0, 100_000_000), step=1_000
    )

    # List all submission folders
    folders = [
        f for f in os.listdir(UPLOAD_FOLDER)
        if os.path.isdir(os.path.join(UPLOAD_FOLDER, f))
    ]

    # Build filtered records for dashboard metrics and charts
    records = []
    for fld in folders:
        base = os.path.join(UPLOAD_FOLDER, fld)

        # Read metadata
        info = {}
        meta_file = os.path.join(base, "info.txt")
        if os.path.exists(meta_file):
            for line in open(meta_file).read().splitlines():
                if ":" in line:
                    k, v = line.split(":", 1)
                    info[k.strip()] = v.strip()

        # Load summary (entrepreneur/admin layered)
        summary_dict, audit = load_layered_summary(base)
        entrepreneur_pending = audit.get("entrepreneur_pending", False)
        # Read current due-diligence status
        status = "Identified"
        status_file = os.path.join(base, "status.json")
        if os.path.exists(status_file):
            try:
                status = json.load(open(status_file)).get("status", status)
            except:
                pass

        # Read status
        status = "Identified"
        spath = os.path.join(base, "status.json")
        if os.path.exists(spath):
            try:
                status = json.load(open(spath)).get("status", status)
            except:
                pass

        # Parse Last 12 months revenues (USD) as numeric or string
        raw_rev = summary_dict.get("Last 12 months revenues (USD)", 0)
        if isinstance(raw_rev, (int, float)):
            rev_val = float(raw_rev)
        else:
            try:
                rev_val = float(str(raw_rev).replace(",", ""))
            except:
                rev_val = 0.0

        # Parse Market size or SOM (USD) as numeric or string
        raw_som = summary_dict.get("Market size or SOM (USD)", 0)
        if isinstance(raw_som, (int, float)):
            som_val = float(raw_som)
        else:
            try:
                som_val = float(str(raw_som).replace(",", ""))
            except:
                som_val = 0.0
        # Parse Expected IRR (%) as numeric or percentage string
        raw_irr = summary_dict.get("Expected IRR (%)", 0)
        if isinstance(raw_irr, (int, float)):
            irr_val = float(raw_irr)
        else:
            try:
                irr_val = float(str(raw_irr).replace("%", "").replace(",", ""))
            except:
                irr_val = 0.0
        try:
            fn_val = float(str(summary_dict.get("Financing need or round size (USD)", "0")).replace(",", ""))
        except:
            fn_val = 0.0

        # Parse portfolio list (may be string, list, or absent)
        port_raw = summary_dict.get("Portfolio", [])
        if isinstance(port_raw, str):
            port_list = [p.strip() for p in port_raw.split(";") if p.strip()]
        elif isinstance(port_raw, list):
            port_list = port_raw
        else:
            port_list = []

        # Parse rejected list (may be string, list, or absent)
        rej_raw = summary_dict.get("Rejected", [])
        if isinstance(rej_raw, str):
            rej_list = [p.strip() for p in rej_raw.split(";") if p.strip()]
        elif isinstance(rej_raw, list):
            rej_list = rej_raw
        else:
            rej_list = []

        # Build filter flags
        flags = []
        # Primary sector / theme
        if primary_sector_filter:
            flags.append(info.get("Sector", "") in primary_sector_filter)
        # Headquarters country
        if hq_country_filter:
            flags.append(info.get("Country HQ", "") in hq_country_filter)
        # Specific Sector(s)
        if specific_sector_filter:
            raw = summary_dict.get("Specific Sector(s)", "")
            if isinstance(raw, list):
                flags.append(any(s in raw for s in specific_sector_filter))
            else:
                flags.append(any(s in str(raw) for s in specific_sector_filter))
        # Region of operation
        if region_filter:
            flags.append(summary_dict.get("Region of operation", "") in region_filter)
        # Main country of current operations
        if main_country_filter:
            flags.append(summary_dict.get("Main country of current operations", "") in main_country_filter)
        # Instrument
        if instrument_filter:
            flags.append(summary_dict.get("Instrument", "") in instrument_filter)
        # Maturity stage
        if maturity_filter:
            flags.append(summary_dict.get("Maturity stage", "") in maturity_filter)
        # Main SDGs targeted (3 max)
        if sdg_filter:
            raw_sdg = summary_dict.get("Main SDGs targeted (3 max)", "")
            sdg_list2 = raw_sdg.split(";") if isinstance(raw_sdg, str) else raw_sdg
            flags.append(any(s in sdg_list2 for s in sdg_filter))
        # Due Diligence stage
        if stage_filter:
            flags.append(status in stage_filter)
        # Numerical filters
        if revenue_min > 0 or revenue_max < 10_000_000:
            flags.append(revenue_min <= rev_val <= revenue_max)
        if som_min > 0 or som_max < 100_000_000:
            flags.append(som_min <= som_val <= som_max)
        if irr_min > 0.0 or irr_max < 100.0:
            flags.append(irr_min <= irr_val <= irr_max)

        # Apply AND/OR logic
        if flags:
            if filter_logic == "AND" and not all(flags):
                continue
            if filter_logic == "OR" and not any(flags):
                continue

        # Last update timestamp
        try:
            adm_path = os.path.join(base, "summary_admin.json")
            ent_path = os.path.join(base, "summary_entrepreneur.json")
            if os.path.exists(adm_path):
                mtime = os.path.getmtime(adm_path)
            elif os.path.exists(ent_path):
                mtime = os.path.getmtime(ent_path)
            else:
                mtime = None
            last_update = datetime.fromtimestamp(mtime) if mtime else None
        except:
            last_update = None

        # Normalize SDG values so downstream code can rely on a list
        raw_sdg = summary_dict.get("Main SDGs targeted (3 max)", "")
        if isinstance(raw_sdg, list):
            sdg_list = raw_sdg
        else:
            sdg_list = [s.strip() for s in str(raw_sdg).split(";") if s.strip()]

        # Append record
        records.append({
            "folder": fld,
            "Project registered name": info.get("Project", fld),
            "Date of incorporation": info.get("Incorporation date", ""),
            "Headquarters country": info.get("Country HQ", ""),
            "Main country of current operations": summary_dict.get("Main country of current operations", ""),
            "Region of operation": summary_dict.get("Region of operation", ""),
            "Primary sector / theme": info.get("Sector", ""),
            "Specific Sector(s)": summary_dict.get("Specific Sector(s)", ""),
            "Impact Area": summary_dict.get("Impact Area", ""),
            "Main SDGs targeted (3 max)": sdg_list,
            "Business Model": summary_dict.get("Business Model", ""),
            "Problem": summary_dict.get("Problem", ""),
            "Solution": summary_dict.get("Solution", ""),
            "Barrier(s) to entry": summary_dict.get("Barrier(s) to entry", ""),
            "Maturity stage": summary_dict.get("Maturity stage", ""),
            "Key risks": summary_dict.get("Key risks", ""),
            "Last 12 months revenues (USD)": rev_val,
            "Breakeven year": summary_dict.get("Breakeven year", ""),
            "Market size or SOM (USD)": som_val,
            "Expected IRR (%)": irr_val,
            "Instrument": summary_dict.get("Instrument", ""),
            "Financing need or round size (USD)": fn_val,
            "Use of proceeds": summary_dict.get("Use of proceeds", ""),
            "Portfolio": "; ".join(port_list),
            "LastUpdate": last_update,
            "Status": status,
            "EntrepreneurEditsPending": entrepreneur_pending,
            "Core team": summary_dict.get("Core team", ""),
            "Contact e‑mail": info.get("Email", "")
        })

    if records:
        df = pd.DataFrame(records)
        # Remove internal folder column (duplicate project identifier)
        df.drop(columns=["folder"], inplace=True)
        # Remove any AI-generated “Project” column – keep only “Project registered name”
        df.drop(columns=["Project"], inplace=True, errors="ignore")
        # Alias "Primary sector / theme" as "Sector" for downstream metrics and charts
        df["Sector"] = df["Primary sector / theme"]
        # =====================================================
        # DASHBOARD SECTION
        # =====================================================
        # Compute extra KPIs
        days_in_stage = []
        for rec in records:
            status_path = os.path.join(UPLOAD_FOLDER, rec["folder"], "status.json")
            if os.path.exists(status_path):
                days_in_stage.append(
                    (datetime.now() - datetime.fromtimestamp(os.path.getmtime(status_path))).days
                )
        avg_days_stage = round(sum(days_in_stage) / len(days_in_stage), 1) if days_in_stage else 0

        financing_col = "Financing need or round size (USD)"
        avg_financing = df[financing_col].mean() if financing_col in df.columns else 0

        sdgs_col = "Main SDGs targeted (3 max)"
        avg_sdgs = df[sdgs_col].apply(lambda x: len(x) if isinstance(x, list) else 0).mean() if sdgs_col in df.columns else 0

        # Year-over-year growth
        df_time = df.dropna(subset=["LastUpdate"]).copy()
        df_time["Year"] = df_time["LastUpdate"].dt.year
        counts = df_time.groupby("Year").size().sort_index()
        if len(counts) >= 2 and counts.iloc[-2] > 0:
            yoy_growth = (counts.iloc[-1] - counts.iloc[-2]) / counts.iloc[-2]
        else:
            yoy_growth = 0

        # KPI Cards Overview
        st.subheader("Traffic Light Indicators")
        df_time["YearMonth"] = df_time["LastUpdate"].dt.to_period("M")
        current_ym = pd.Period(datetime.now(), "M")
        prev_ym = current_ym - 1
        curr_subs = df_time[df_time["YearMonth"] == current_ym].shape[0]
        prev_subs = df_time[df_time["YearMonth"] == prev_ym].shape[0]
        delta_subs = curr_subs - prev_subs

        def _mean_safe(series): return series.mean() if len(series) else 0
        revenue_col = "Last 12 months revenues (USD)"
        irr_col = "Expected IRR (%)"
        if revenue_col in df_time.columns:
            rev_curr = _mean_safe(df_time[df_time["YearMonth"] == current_ym][revenue_col])
            rev_prev = _mean_safe(df_time[df_time["YearMonth"] == prev_ym][revenue_col])
        else:
            rev_curr = rev_prev = 0.0
        # Use the "Expected IRR (%)" column for IRR metrics
        if irr_col in df_time.columns:
            irr_curr = _mean_safe(df_time[df_time["YearMonth"] == current_ym][irr_col])
            irr_prev = _mean_safe(df_time[df_time["YearMonth"] == prev_ym][irr_col])
        else:
            irr_curr = irr_prev = 0.0

        c1, c2, c3, c4 = st.columns(4)
        status_emoji = "🟢" if delta_subs > 0 else ("⚠️" if delta_subs == 0 else "🔴")
        c1.metric(f"{status_emoji} Total Projects", len(df), f"{delta_subs:+d} MoM")
        rev_delta = rev_curr - rev_prev
        c2.metric("💰 Avg Revenue (USD)", f"${df[revenue_col].mean():,.0f}", f"{rev_delta:+,.0f}")
        irr_med = df[irr_col].median() if irr_col in df.columns else 0
        irr_delta = irr_curr - irr_prev
        c3.metric("📈 Median IRR (%)", f"{irr_med:.1f}%", f"{irr_delta:+.1f}%")
        late_pct = (sum(d > 30 for d in days_in_stage) / len(days_in_stage) * 100) if days_in_stage else 0
        color_late = "🔴" if late_pct > 20 else ("⚠️" if late_pct > 10 else "🟢")
        c4.metric(f"{color_late} % Pipeline Late", f"{late_pct:.1f}%")

        k1, k2, k3 = st.columns(3)
        k1.metric("🌐 Unique Sectors", df["Sector"].nunique())
        k2.metric("💼 Avg Financing (USD)", f"${avg_financing:,.0f}")
        k3.metric("📊 Avg SDGs / Project", f"{avg_sdgs:.1f}")

        # At-a-glance visuals – compact 2×2 grid
        st.subheader("At a Glance")
        # First row: status and regional distribution
        r1c1, r1c2 = st.columns(2)
        with r1c1:
            st.markdown("**Due Diligence Status**")
            status_counts = df["Status"].value_counts() if "Status" in df.columns else pd.Series(dtype=int)
            fig1 = plt.figure(figsize=(2.5, 2.5))
            plt.pie(status_counts, labels=status_counts.index, autopct="%1.0f%%",
                    colors=[STATUS_COLORS.get(s, "#cccccc") for s in status_counts.index])
            st.pyplot(fig1, use_container_width=True)
        with r1c2:
            st.markdown("**Regional Distribution**")
            region_counts = df["Region of operation"].value_counts() if "Region of operation" in df.columns else pd.Series(dtype=int)
            fig2 = plt.figure(figsize=(2.5, 2.5))
            plt.pie(region_counts, labels=region_counts.index, autopct="%1.0f%%",
                    colors=[REGION_COLORS.get(r, "#cccccc") for r in region_counts.index])
            st.pyplot(fig2, use_container_width=True)
        # Second row: sector submissions and pipeline funnel
        r2c1, r2c2 = st.columns(2)
        with r2c1:
            st.markdown("**Submissions per Sector**")
            sector_counts = df["Sector"].value_counts() if "Sector" in df.columns else pd.Series(dtype=int)
            fig3 = plt.figure(figsize=(4, 2.5))
            plt.bar(sector_counts.index, sector_counts.values,
                    color=[SECTOR_COLORS.get(s, "#cccccc") for s in sector_counts.index])
            plt.xticks(rotation=45, ha="right")
            st.pyplot(fig3, use_container_width=True)
        with r2c2:
            st.markdown("**Pipeline Funnel**")
            stage_order = REVIEW_STAGES
            funnel_counts = df["Status"].value_counts().reindex(stage_order).fillna(0) if "Status" in df.columns else pd.Series(0, index=stage_order)
            fig4 = plt.figure(figsize=(3.5, 2.5))
            plt.barh(funnel_counts.index[::-1], funnel_counts.values[::-1],
                     color=[STATUS_COLORS.get(s, "#cccccc") for s in funnel_counts.index[::-1]])
            plt.xlabel("Projects")
            plt.tight_layout()
            st.pyplot(fig4, use_container_width=True)

        # Projects table
        st.subheader("Projects Table")
        # Convert list‑type cells to strings so PyArrow can serialize the dataframe
        df_display = df.applymap(
            lambda v: "; ".join(map(str, v)) if isinstance(v, list) else v
        )
        # Remove duplicate default columns from the display
        df_display.drop(
            columns=["Revenues", "SOM", "IRR", "FinancingNeed", "Trend", "Rejected", "Sector"],
            inplace=True,
            errors="ignore"
        )
        all_cols = df_display.columns.tolist()
        styled = df_display[all_cols].style \
            .applymap(style_status, subset=[c for c in ["Status"] if c in all_cols]) \
            .applymap(style_region, subset=[c for c in ["Region of operation"] if c in all_cols]) \
            .applymap(style_sector, subset=[c for c in ["Sector"] if c in all_cols])
        st.dataframe(styled, use_container_width=True)

        # ── Projects Table ──────────────────────────────────────────
        st.subheader("Projects Table")
        styled = df[[
            "Project registered name", "Sector", "Region of operation", "Status",
            "Portfolio", "Last 12 months revenues (USD)", "Market size or SOM (USD)",
            "Expected IRR (%)", "Financing need or round size (USD)", "LastUpdate"
        ]].style \
            .applymap(style_status, subset=["Status"]) \
            .applymap(style_region, subset=["Region of operation"]) \
            .applymap(style_sector, subset=["Sector"])
        st.dataframe(styled, use_container_width=True)

        # ── Detailed “View & Edit Projects” interface ──────────────────
        for row_num, rec in enumerate(records):
            fld = rec["folder"]
            base = os.path.join(UPLOAD_FOLDER, fld)

            # Re‑read metadata for safety
            info = {}
            meta_file = os.path.join(base, "info.txt")
            if os.path.exists(meta_file):
                for line in open(meta_file).read().splitlines():
                    if ":" in line:
                        k, v = line.split(":", 1)
                        info[k.strip()] = v.strip()

            # Layered summary (entrepreneur / admin)
            summary_dict, audit = load_layered_summary(base)

            entrepreneur_pending = rec.get("EntrepreneurEditsPending")
            # Expander header with running index + project name + last update + pending indicator
            lu = rec.get("LastUpdate")
            lu_str = lu.strftime("%Y-%m-%d %H:%M") if pd.notnull(lu) else "Unknown"
            indicator = " ⚠️" if rec.get("EntrepreneurEditsPending") else ""
            exp_title = f"{row_num}. {rec['Project registered name']}{indicator} — Last update: {lu_str}"
            with st.expander(exp_title, expanded=False):
                # ── Stage selector (auto‑save) ──────────────────────────
                status_path = os.path.join(base, "status.json")
                cur_stage = rec.get("Status", "Identified")
                try:
                    idx = REVIEW_STAGES.index(cur_stage)
                except ValueError:
                    idx = 0
                new_stage = st.selectbox(
                    "Due Diligence / Operations stage",
                    REVIEW_STAGES,
                    index=idx,
                    key=f"stage_{fld}",
                    help="Changes are saved immediately",
                )
                if new_stage != cur_stage:
                    with open(status_path, "w", encoding="utf-8") as sf:
                        json.dump({"status": new_stage}, sf)
                    st.success("Status updated!")
                    _rerun()

                # ── Summary grid (read‑only) ───────────────────────────
                st.markdown("#### Project Summary")
                # Merge submission metadata with AI summary for display
                meta_display = {
                    "Project registered name": info.get("Project", ""),
                    "Date of incorporation": info.get("Incorporation date", ""),
                    "Primary sector / theme": info.get("Sector", ""),
                    "Headquarters country": info.get("Country HQ", ""),
                    "Contact e-mail": info.get("Email", ""),
                }
                combined = {**meta_display, **summary_dict}
                render_summary_grid(combined)

                show_versions = st.checkbox("Show version history", key=f"show_versions_{fld}")
                if show_versions:
                    # ---------- Version tabs (sorted by timestamp, newest left; GPT original last) ----------
                    entries = []
                    admin_path = os.path.join(base, "summary_admin.json")
                    ent_path   = os.path.join(base, "summary_entrepreneur.json")
                    gpt_path   = os.path.join(base, "summary_gpt.txt")
                    if os.path.exists(admin_path):
                        adm = json.load(open(admin_path, encoding="utf-8"))
                        # Admin current
                        try:
                            ts_cur = datetime.fromisoformat(adm.get("audit", {}).get("admin_last_updated")).timestamp()
                        except:
                            ts_cur = 0
                        entries.append({"label": "Admin (current)", "data": adm.get("fields", {}), "ts": ts_cur})
                        # Admin history
                        for idx, item in enumerate(reversed(adm.get("history", [])), 1):
                            try:
                                ts_hist = datetime.fromisoformat(item.get("audit", {}).get("admin_last_updated")).timestamp()
                            except:
                                ts_hist = 0
                            dt_str = item.get("audit", {}).get("admin_last_updated", "")[:10]
                            entries.append({"label": f"Admin v{idx} {dt_str}", "data": item.get("fields", {}), "ts": ts_hist})
                    if os.path.exists(ent_path):
                        ent = json.load(open(ent_path, encoding="utf-8"))
                        # Entrepreneur current
                        try:
                            ts_cur = datetime.fromisoformat(ent.get("audit", {}).get("entrepreneur_last_updated")).timestamp()
                        except:
                            ts_cur = 0
                        entries.append({"label": "Entrepreneur (current)", "data": ent.get("fields", {}), "ts": ts_cur})
                        # Entrepreneur history
                        for idx, item in enumerate(reversed(ent.get("history", [])), 1):
                            try:
                                ts_hist = datetime.fromisoformat(item.get("audit", {}).get("entrepreneur_last_updated")).timestamp()
                            except:
                                ts_hist = 0
                            dt_str = item.get("audit", {}).get("entrepreneur_last_updated", "")[:10]
                            entries.append({"label": f"Ent v{idx} {dt_str}", "data": item.get("fields", {}), "ts": ts_hist})
                    if os.path.exists(gpt_path):
                        gpt_data = _parse_json_from_string(open(gpt_path, encoding="utf-8").read())
                        entries.append({"label": "GPT original", "data": gpt_data, "ts": 0})
                    # Sort by timestamp descending; GPT original (ts=0) ends up last
                    entries_sorted = sorted(entries, key=lambda e: e["ts"], reverse=True)
                    # Render tabs in order
                    tabs_labels = [e["label"] for e in entries_sorted]
                    tabs_objs = st.tabs(tabs_labels)
                    for e, tab in zip(entries_sorted, tabs_objs):
                        with tab:
                            render_summary_grid(e["data"])

                # ── Download all uploaded documents ────────────────────
                st.markdown("#### Uploaded Documents")
                all_docs = [
                    d for d in os.listdir(base)
                    if d.lower().endswith(('.pdf', '.docx', '.pptx', '.xlsx'))
                ]
                if all_docs:
                    for doc in all_docs:
                        with open(os.path.join(base, doc), "rb") as dbin:
                            st.download_button(
                                label=f"⬇️ {doc}",
                                data=dbin.read(),
                                file_name=doc,
                                mime="application/octet-stream",
                                key=f"dl_{fld}_{doc}",
                            )
                else:
                    st.info("No documents uploaded.")

                # ── Admin edit form ────────────────────────────────────
                st.markdown("#### Edit Fields")
                with st.form(f"edit_{fld}"):
                    # Edit submission metadata
                    meta_proj = st.text_input(
                        "Project registered name",
                        value=info.get("Project", ""),
                        key=f"meta_proj_{fld}"
                    )
                    meta_date = st.date_input(
                        "Date of incorporation",
                        value=datetime.fromisoformat(info.get("Incorporation date", datetime.now().date().isoformat())),
                        key=f"meta_date_{fld}"
                    )
                    meta_sector = st.selectbox(
                        "Primary sector / theme",
                        SECTOR_OPTIONS,
                        index=SECTOR_OPTIONS.index(info.get("Sector", "")) if info.get("Sector", "") in SECTOR_OPTIONS else 0,
                        key=f"meta_sector_{fld}"
                    )
                    meta_country = st.selectbox(
                        "Headquarters country",
                        COUNTRY_OPTIONS,
                        index=COUNTRY_OPTIONS.index(info.get("Country HQ", "")) if info.get("Country HQ", "") in COUNTRY_OPTIONS else 0,
                        key=f"meta_country_{fld}"
                    )
                    meta_email = st.text_input(
                        "Contact e‑mail",
                        value=info.get("Email", ""),
                        key=f"meta_email_{fld}"
                    )
                    # Update metadata dict
                    info["Project"] = meta_proj
                    info["Incorporation date"] = meta_date.isoformat()
                    info["Sector"] = meta_sector
                    info["Country HQ"] = meta_country
                    info["Email"] = meta_email
                    edited = {}
                    for field in AI_FIELDS:
                        val = summary_dict.get(field, "")
                        if field == "Maturity stage":
                            edited[field] = st.selectbox(
                                field,
                                MATURITY_STAGES,
                                index=MATURITY_STAGES.index(val) if val in MATURITY_STAGES else 0,
                            )
                        elif field == "Main SDGs targeted (3 max)":
                            defaults = (
                                [s.strip() for s in val.split(";")] if isinstance(val, str) else val
                            )
                            safe_defaults = [s for s in defaults if s in SDG_OPTIONS]
                            edited[field] = st.multiselect(
                                field, SDG_OPTIONS, default=safe_defaults, max_selections=3
                            )
                        elif field in [
                            "Last 12 months revenues (USD)",
                            "Market size or SOM (USD)",
                            "Financing need or round size (USD)",
                        ]:
                            try:
                                num_val = float(str(val).replace(",", ""))
                            except Exception:
                                num_val = 0.0
                            edited[field] = st.number_input(
                                field, min_value=0.0, value=num_val, step=1000.0
                            )
                        elif field == "Expected IRR (%)":
                            try:
                                num_val = float(str(val).replace("%", ""))
                            except Exception:
                                num_val = 0.0
                            edited[field] = st.number_input(
                                field, min_value=0.0, max_value=100.0, value=num_val, step=0.1
                            )
                        elif field == "Breakeven year":
                            # Parse existing value or default to current year
                            try:
                                year_val = int(val)
                            except:
                                year_val = datetime.now().year
                            # Clamp to min/max bounds
                            year_val = max(1900, min(year_val, 2100))
                            edited[field] = st.number_input(
                                "Breakeven year",
                                min_value=1900,
                                max_value=2100,
                                value=year_val,
                                step=1,
                                format="%d"
                            )
                        elif field == "Region of operation":
                            edited[field] = st.selectbox(
                                field,
                                list(REGION_COLORS.keys()),
                                index=list(REGION_COLORS.keys()).index(val) if val in REGION_COLORS else 0,
                            )
                        elif field == "Main country of current operations":
                            edited[field] = st.selectbox(
                                field,
                                COUNTRY_OPTIONS,
                                index=COUNTRY_OPTIONS.index(val) if val in COUNTRY_OPTIONS else 0,
                            )
                        elif field == "Instrument":
                            edited[field] = st.selectbox(
                                field,
                                INSTRUMENT_OPTIONS,
                                index=INSTRUMENT_OPTIONS.index(val) if val in INSTRUMENT_OPTIONS else 0,
                            )
                        elif field == "Specific Sector(s)":
                            defaults = ([s.strip() for s in val.split(";")]
                                        if isinstance(val, str) else val)
                            safe_defaults = [s for s in defaults if s in SPECIFIC_SECTOR_OPTIONS]
                            edited[field] = st.multiselect(
                                field, SPECIFIC_SECTOR_OPTIONS, default=safe_defaults
                            )
                        elif field == "Use of proceeds":
                            edited[field] = st.text_area(field, value=str(val))
                        else:
                            edited[field] = st.text_input(field, value=str(val))

                    if st.form_submit_button("💾 Save edits"):
                        # Save updated metadata
                        meta_file = os.path.join(base, "info.txt")
                        meta_text = "\n".join(f"{k}: {v}" for k, v in info.items()) + "\nNDA: Accepted\n"
                        with open(meta_file, "w", encoding="utf-8") as mf:
                            mf.write(meta_text)
                        # Save updated admin summary
                        _save_admin_summary(fld, edited)
                        st.success("Edits saved!")
                        _rerun()
# ═════════════════════════════════════════════════════════════════
# PORTFOLIO‑HOLDER DASHBOARD (NCGE/NCGD)
# ═════════════════════════════════════════════════════════════════
elif is_portfolio_ncge or is_portfolio_ncgd:
    if is_portfolio_ncge:
        st.title("Nature Catalyst Growth :orange[Equity] Portfolio")
    elif is_portfolio_ncgd:
        st.title("Nature Catalyst Growth :orange[Debt] Portfolio")
    # Determine the role tag
    role_tag = "NCGE" if is_portfolio_ncge else "NCGD"

    # List all submission folders
    folders = [
        f for f in os.listdir(UPLOAD_FOLDER)
        if os.path.isdir(os.path.join(UPLOAD_FOLDER, f))
    ]

    records = []
    for fld in folders:
        base = os.path.join(UPLOAD_FOLDER, fld)
        info = {}
        meta_file = os.path.join(base, "info.txt")
        if os.path.exists(meta_file):
            for line in open(meta_file).read().splitlines():
                if ":" in line:
                    k, v = line.split(":", 1)
                    info[k.strip()] = v.strip()
        summary_dict, audit = load_layered_summary(base)
        entrepreneur_pending = audit.get("entrepreneur_pending", False)
        # Read current due-diligence status
        status = "Identified"
        status_path = os.path.join(base, "status.json")
        if os.path.exists(status_path):
            try:
                status = json.load(open(status_path)).get("status", status)
            except:
                pass
        # Parse portfolio list
        port_raw = summary_dict.get("Portfolio", [])
        if isinstance(port_raw, str):
            portfolio_list = [p.strip() for p in port_raw.split(";") if p.strip()]
        elif isinstance(port_raw, list):
            portfolio_list = port_raw
        else:
            portfolio_list = []
        # Rejected list
        rej_raw = summary_dict.get("Rejected", [])
        if isinstance(rej_raw, str):
            rejected_list = [p.strip() for p in rej_raw.split(";") if p.strip()]
        elif isinstance(rej_raw, list):
            rejected_list = rej_raw
        else:
            rejected_list = []
        # SDGs
        raw_sdg = summary_dict.get("Main SDGs targeted (3 max)", "")
        if isinstance(raw_sdg, list):
            sdg_list = raw_sdg
        else:
            sdg_list = [s.strip() for s in str(raw_sdg).split(";") if s.strip()]
        # Determine last update timestamp for portfolio entry
        try:
            adm_path = os.path.join(base, "summary_admin.json")
            ent_path = os.path.join(base, "summary_entrepreneur.json")
            if os.path.exists(adm_path):
                mtime = os.path.getmtime(adm_path)
            elif os.path.exists(ent_path):
                mtime = os.path.getmtime(ent_path)
            else:
                mtime = None
            last_update = datetime.fromtimestamp(mtime) if mtime else None
        except:
            last_update = None
        # Parse numeric fields for portfolio entries
        raw_rev = summary_dict.get("Last 12 months revenues (USD)", 0)
        if isinstance(raw_rev, (int, float)):
            rev_val = float(raw_rev)
        else:
            try:
                rev_val = float(str(raw_rev).replace(",", ""))
            except:
                rev_val = 0.0

        raw_som = summary_dict.get("Market size or SOM (USD)", 0)
        if isinstance(raw_som, (int, float)):
            som_val = float(raw_som)
        else:
            try:
                som_val = float(str(raw_som).replace(",", ""))
            except:
                som_val = 0.0

        raw_irr = summary_dict.get("Expected IRR (%)", 0)
        if isinstance(raw_irr, (int, float)):
            irr_val = float(raw_irr)
        else:
            try:
                irr_val = float(str(raw_irr).replace("%", "").replace(",", ""))
            except:
                irr_val = 0.0
        # Append
        records.append({
            "folder": fld,
            "Project registered name": info.get("Project", fld),
            "Date of incorporation": info.get("Incorporation date", ""),
            "Primary sector / theme": info.get("Sector", ""),
            "Sector": info.get("Sector", ""),
            "Headquarters country": info.get("Country HQ", ""),
            "Contact e‑mail": info.get("Email", ""),
            "Specific Sector(s)": summary_dict.get("Specific Sector(s)", ""),
            "Main country of current operations": summary_dict.get("Main country of current operations", ""),
            "Business Model": summary_dict.get("Business Model", ""),
            "Core team": summary_dict.get("Core team", ""),
            "Key risks": summary_dict.get("Key risks", ""),
            "Last 12 months revenues (USD)": rev_val,
            "Breakeven year": summary_dict.get("Breakeven year", ""),
            "Market size or SOM (USD)": som_val,
            "Expected IRR (%)": irr_val,
            "Financing need or round size (USD)": summary_dict.get("Financing need or round size (USD)", ""),
            "Instrument": summary_dict.get("Instrument", ""),
            "Use of proceeds": summary_dict.get("Use of proceeds", ""),
            "Impact Area": summary_dict.get("Impact Area", ""),
            "Main SDGs targeted (3 max)": sdg_list,
            "Solution": summary_dict.get("Solution", ""),
            "Barrier(s) to entry": summary_dict.get("Barrier(s) to entry", ""),
            "Region of operation": summary_dict.get("Region of operation", ""),
            "Status": status,
            "Maturity stage": summary_dict.get("Maturity stage", ""),
            "Portfolio": portfolio_list,
            "Rejected": rejected_list,
            "LastUpdate": last_update,
            "EntrepreneurEditsPending": entrepreneur_pending,
        })
    # Re-classify projects
    unclassified = [
        r for r in records
        if role_tag not in r["Portfolio"] and role_tag not in r["Rejected"]
    ]
    in_portfolio = [r for r in records if role_tag in r["Portfolio"]]
    rejected = [r for r in records if role_tag in r["Rejected"]]

    # Tabs: Unclassified, Portfolio, Rejected
    tabs = st.tabs(["🆕 Unclassified", "📂 Portfolio", "🚫 Rejected"])

    # -------------------- Tab 1: Unclassified -----------------------
    with tabs[0]:
        st.subheader(f"Unclassified (_{len(unclassified)} projects_)", divider="gray")
        st.caption("➕ **Add** → portfolio · ❌ **Eject** → Rejected (can reconsider later).")
        if not unclassified:
            st.info("No unclassified projects.")
        else:
            for rec in unclassified:
                with st.expander(rec["Project registered name"]):
                    base_path = os.path.join(UPLOAD_FOLDER, rec["folder"])
                    summ, _ = load_layered_summary(base_path)
                    # Merge metadata
                    meta_display = {
                        "Project registered name": rec.get("Project registered name", rec.get("Project Name", "")),
                        "Date of incorporation": rec.get("Date of incorporation", ""),
                        "Primary sector / theme": rec.get("Primary sector / theme", ""),
                        "Headquarters country": rec.get("Headquarters country", ""),
                        "Contact e-mail": rec.get("Contact e-mail", ""),
                    }
                    combined = {**meta_display, **summ}
                    render_summary_grid(combined)
                    colA, colB = st.columns([1,1])
                    add_key = f"add_{role_tag}_{rec['folder']}"
                    eject_key = f"eject_{role_tag}_{rec['folder']}"
                    if colA.button("➕ Add to Portfolio", key=add_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "add")
                        st.success("Added to portfolio.")
                        _rerun()
                    if colB.button("❌ Eject", key=eject_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "eject")
                        st.warning("Ejected.")
                        _rerun()

    # -------------------- Tab 2: My Portfolio -----------------------
    with tabs[1]:
        # --- BEGIN: Filter Sidebar (read-only) ---
        # Sidebar filters for portfolio holder
        with st.sidebar:
            st.markdown("### 🔍 Filter Projects")
            logic = st.radio("Filter logic", ["AND", "OR"], horizontal=True)
            st.sidebar.markdown("#### Categorical Filters")
            
            primary_sector_filter = st.multiselect(
                "Primary sector / theme", SECTOR_OPTIONS
            )
            
            specific_sector_filter = st.multiselect(
                "Specific Sector(s)", SPECIFIC_SECTOR_OPTIONS
            )

            hq_country_filter = st.multiselect(
                "Headquarters country", COUNTRY_OPTIONS
            )
            
            region_filter = st.multiselect(
                "Region of operation", list(REGION_COLORS.keys())
            )
            main_country_filter = st.multiselect(
                "Main country of current operations", COUNTRY_OPTIONS
            )
            instrument_filter = st.multiselect(
                "Instrument", INSTRUMENT_OPTIONS
            )
            maturity_filter = st.multiselect(
                "Maturity stage", MATURITY_STAGES
            )
            sdg_filter = st.multiselect(
                "Main SDGs targeted (3 max)", SDG_OPTIONS
            )
            stage_filter = st.multiselect(
                "Due Diligence stage", REVIEW_STAGES
            )
            st.sidebar.markdown("#### Numerical Filters")
            irr_min, irr_max = st.slider(
                "IRR Range (%)", min_value=0.0, max_value=100.0, value=(0.0, 100.0), step=0.1
            )
            revenue_min, revenue_max = st.slider(
                "Revenue Range (USD)", min_value=0, max_value=1_000_000_000, value=(0, 1_000_000_000), step=1_000
            )
            som_min, som_max = st.slider(
                "SOM Range (USD)", min_value=0, max_value=1_000_000_000, value=(0, 1_000_000_000), step=1_000
            )
        # --- END: Filter Sidebar ---

        # --- BEGIN: Apply filters to in_portfolio ---
        def portfolio_filter_logic(rec):
            flags = []
            # Primary sector / theme
            if primary_sector_filter:
                flags.append(rec.get("Sector", "") in primary_sector_filter)
            # Specific Sector(s)
            if specific_sector_filter:
                raw = rec.get("Specific Sector(s)", "")
                if isinstance(raw, list):
                    flags.append(any(s in raw for s in specific_sector_filter))
                else:
                    flags.append(any(s in str(raw) for s in specific_sector_filter))
            # Headquarters country
            if hq_country_filter:
                flags.append(rec.get("Headquarters country", "") in hq_country_filter)
            # Region of operation
            if region_filter:
                flags.append(rec.get("Region of operation", "") in region_filter)
            # Main country of current operations
            if main_country_filter:
                flags.append(rec.get("Main country of current operations", "") in main_country_filter)
            # Instrument
            if instrument_filter:
                flags.append(rec.get("Instrument", "") in instrument_filter)
            # Maturity stage
            if maturity_filter:
                flags.append(rec.get("Maturity stage", "") in maturity_filter)
            # Main SDGs targeted (3 max)
            if sdg_filter:
                raw_sdg = rec.get("Main SDGs targeted (3 max)", "")
                sdg_list2 = raw_sdg.split(";") if isinstance(raw_sdg, str) else raw_sdg
                flags.append(any(s in sdg_list2 for s in sdg_filter))
            # Due Diligence stage
            if stage_filter:
                flags.append(rec.get("Status", "") in stage_filter)
            # Numerical filters
            try:
                irr_val = float(str(rec.get("Expected IRR (%)", 0)).replace("%", "").replace(",", ""))
            except:
                irr_val = 0.0
            try:
                rev_val = float(str(rec.get("Last 12 months revenues (USD)", 0)).replace(",", ""))
            except:
                rev_val = 0.0
            try:
                som_val = float(str(rec.get("Market size or SOM (USD)", 0)).replace(",", ""))
            except:
                som_val = 0.0
            flags.append(irr_min <= irr_val <= irr_max)
            flags.append(revenue_min <= rev_val <= revenue_max)
            flags.append(som_min <= som_val <= som_max)
            # Apply filter logic
            if not flags:
                return True
            if logic == "AND":
                return all(flags)
            else:
                return any(flags)

        filtered_portfolio = [rec for rec in in_portfolio if portfolio_filter_logic(rec)]
        # --- END: Apply filters

        st.subheader(f"Portfolio (_{len(filtered_portfolio)} projects_)", divider="gray")
        if not filtered_portfolio:
            st.info("No projects in your portfolio match the current filters.")
        else:
            dfp = pd.DataFrame(filtered_portfolio)
            # Remove internal folder column from portfolio table
            dfp.drop(columns=["folder"], inplace=True)
            # Remove any AI-generated “Project” column – keep only “Project registered name”
            dfp.drop(columns=["Project"], inplace=True, errors="ignore")
            # Alias "Primary sector / theme" as "Sector" for styling consistency
            dfp["Sector"] = dfp["Primary sector / theme"]

            # Projects Table (Portfolio) — mirror Admin layout
            dfp_display = dfp.applymap(
                lambda v: "; ".join(map(str, v)) if isinstance(v, list) else v
            )
            dfp_display.drop(
                columns=["Revenues", "SOM", "IRR", "FinancingNeed", "Trend", "Rejected"],
                inplace=True,
                errors="ignore"
            )
            desired_order = [
                "Project registered name","Date of incorporation","Headquarters country",
                "Main country of current operations","Region of operation","Primary sector / theme",
                "Specific Sector(s)","Impact Area","Main SDGs targeted (3 max)",
                "Business Model","Problem","Solution","Barrier(s) to entry",
                "Maturity stage","Key risks","Last 12 months revenues (USD)",
                "Breakeven year","Market size or SOM (USD)","Expected IRR (%)",
                "Instrument","Financing need or round size (USD)","Use of proceeds",
                "Portfolio","LastUpdate","Status","Core team","Contact e-mail"
            ]

            # Projects table (Portfolio) under Traffic Light Indicators
            st.subheader("Projects Table")
            # Convert list‑type cells to strings for display
            dfp_display = dfp.applymap(
                lambda v: "; ".join(map(str, v)) if isinstance(v, list) else v
            )
            # Remove duplicate default columns
            dfp_display.drop(
                columns=["Revenues", "SOM", "IRR", "FinancingNeed", "Trend", "Rejected", "Sector"],
                inplace=True,
                errors="ignore"
            )
            all_cols = dfp_display.columns.tolist()
            styled = dfp_display[all_cols].style \
                .applymap(style_status, subset=[c for c in ["Status"] if c in all_cols]) \
                .applymap(style_region, subset=[c for c in ["Region of operation"] if c in all_cols]) \
                .applymap(style_sector, subset=[c for c in ["Sector"] if c in all_cols])
            st.dataframe(styled, use_container_width=True)

            # Compute extra KPIs using only portfolio projects
            days_in_stage = []
            for rec in filtered_portfolio:
                status_path = os.path.join(UPLOAD_FOLDER, rec["folder"], "status.json")
                if os.path.exists(status_path):
                    days_in_stage.append(
                        (datetime.now() - datetime.fromtimestamp(os.path.getmtime(status_path))).days
                    )
            avg_days_stage = round(sum(days_in_stage) / len(days_in_stage), 1) if days_in_stage else 0

            financing_col = "Financing need or round size (USD)"
            avg_financing = dfp[financing_col].mean() if financing_col in dfp.columns else 0

            sdgs_col = "Main SDGs targeted (3 max)"
            avg_sdgs = dfp[sdgs_col].apply(lambda x: len(x) if isinstance(x, list) else 0).mean() if sdgs_col in dfp.columns else 0

            # Year-over-year growth of portfolio submissions
            dfp_time = dfp.dropna(subset=["LastUpdate"]).copy()
            dfp_time["Year"] = dfp_time["LastUpdate"].dt.year
            # Ensure numeric fields are coerced to numbers to avoid mixing types
            dfp_time["Expected IRR (%)"] = pd.to_numeric(dfp_time["Expected IRR (%)"], errors="coerce")
            dfp_time["Last 12 months revenues (USD)"] = pd.to_numeric(dfp_time["Last 12 months revenues (USD)"], errors="coerce")
            dfp_time["Market size or SOM (USD)"] = pd.to_numeric(dfp_time["Market size or SOM (USD)"], errors="coerce")
            counts = dfp_time.groupby("Year").size().sort_index()
            if len(counts) >= 2 and counts.iloc[-2] > 0:
                yoy_growth = (counts.iloc[-1] - counts.iloc[-2]) / counts.iloc[-2]
            else:
                yoy_growth = 0

            # KPI Cards Overview
            st.subheader("Traffic Light Indicators")
            dfp_time["YearMonth"] = dfp_time["LastUpdate"].dt.to_period("M")
            current_ym = pd.Period(datetime.now(), "M")
            prev_ym = current_ym - 1
            curr_subs = dfp_time[dfp_time["YearMonth"] == current_ym].shape[0]
            prev_subs = dfp_time[dfp_time["YearMonth"] == prev_ym].shape[0]
            delta_subs = curr_subs - prev_subs

            def _mean_safe(series): return series.mean() if len(series) else 0
            revenue_col = "Last 12 months revenues (USD)"
            irr_col = "Expected IRR (%)"
            if revenue_col in dfp_time.columns:
                rev_curr = _mean_safe(dfp_time[dfp_time["YearMonth"] == current_ym][revenue_col])
                rev_prev = _mean_safe(dfp_time[dfp_time["YearMonth"] == prev_ym][revenue_col])
            else:
                rev_curr = rev_prev = 0.0
            if irr_col in dfp_time.columns:
                irr_curr = _mean_safe(dfp_time[dfp_time["YearMonth"] == current_ym][irr_col])
                irr_prev = _mean_safe(dfp_time[dfp_time["YearMonth"] == prev_ym][irr_col])
            else:
                irr_curr = irr_prev = 0.0

            c1, c2, c3, c4 = st.columns(4)
            status_emoji = "🟢" if delta_subs > 0 else ("⚠️" if delta_subs == 0 else "🔴")
            c1.metric(f"{status_emoji} Total Projects", len(dfp), f"{delta_subs:+d} MoM")
            rev_delta = rev_curr - rev_prev
            c2.metric("💰 Avg Revenue (USD)", f"${dfp[revenue_col].mean():,.0f}", f"{rev_delta:+,.0f}")
            irr_med = dfp[irr_col].median() if irr_col in dfp.columns else 0
            irr_delta = irr_curr - irr_prev
            c3.metric("📈 Median IRR (%)", f"{irr_med:.1f}%", f"{irr_delta:+.1f}%")
            late_pct = (sum(d > 30 for d in days_in_stage) / len(days_in_stage) * 100) if days_in_stage else 0
            color_late = "🔴" if late_pct > 20 else ("⚠️" if late_pct > 10 else "🟢")
            c4.metric(f"{color_late} % Pipeline Late", f"{late_pct:.1f}%")

            k1, k2, k3 = st.columns(3)
            k1.metric("🌐 Unique Sectors", dfp["Sector"].nunique())
            k2.metric("💼 Avg Financing (USD)", f"{avg_financing:,.0f}")
            k3.metric("📊 Avg SDGs / Project", f"{avg_sdgs:.1f}")

            # At-a-glance visuals – compact 2×2 grid
            st.subheader("At a Glance")
            # First row: status and regional distribution
            r1c1, r1c2 = st.columns(2)
            with r1c1:
                st.markdown("**Due Diligence Status**")
                status_counts = dfp["Status"].value_counts() if "Status" in dfp.columns else pd.Series(dtype=int)
                fig1 = plt.figure(figsize=(2.5, 2.5))
                plt.pie(status_counts, labels=status_counts.index, autopct="%1.0f%%",
                        colors=[STATUS_COLORS.get(s, "#cccccc") for s in status_counts.index])
                st.pyplot(fig1, use_container_width=True)
            with r1c2:
                st.markdown("**Regional Distribution**")
                region_counts = dfp["Region of operation"].value_counts() if "Region of operation" in dfp.columns else pd.Series(dtype=int)
                fig2 = plt.figure(figsize=(2.5, 2.5))
                plt.pie(region_counts, labels=region_counts.index, autopct="%1.0f%%",
                        colors=[REGION_COLORS.get(r, "#cccccc") for r in region_counts.index])
                st.pyplot(fig2, use_container_width=True)
            # Second row: sector submissions and pipeline funnel
            r2c1, r2c2 = st.columns(2)
            with r2c1:
                st.markdown("**Submissions per Sector**")
                sector_counts = dfp["Sector"].value_counts() if "Sector" in dfp.columns else pd.Series(dtype=int)
                fig3 = plt.figure(figsize=(4, 2.5))
                plt.bar(sector_counts.index, sector_counts.values,
                        color=[SECTOR_COLORS.get(s, "#cccccc") for s in sector_counts.index])
                plt.xticks(rotation=45, ha="right")
                st.pyplot(fig3, use_container_width=True)
            with r2c2:
                st.markdown("**Pipeline Funnel**")
                stage_order = REVIEW_STAGES
                funnel_counts = dfp["Status"].value_counts().reindex(stage_order).fillna(0) if "Status" in dfp.columns else pd.Series(0, index=stage_order)
                fig4 = plt.figure(figsize=(3.5, 2.5))
                plt.barh(funnel_counts.index[::-1], funnel_counts.values[::-1],
                         color=[STATUS_COLORS.get(s, "#cccccc") for s in funnel_counts.index[::-1]])
                plt.xlabel("Projects")
                plt.tight_layout()
                st.pyplot(fig4, use_container_width=True)

            # ── Projects Table (Portfolio) ─────────────────────────────────
            st.subheader("Projects Table")
            # Convert list-type cells to strings for display
            dfp_display = dfp.applymap(
                lambda v: "; ".join(map(str, v)) if isinstance(v, list) else v
            )
            # Drop duplicate default columns (same as Admin)
            dfp_display.drop(
                columns=["Revenues", "SOM", "IRR", "FinancingNeed", "Trend", "Rejected", "Sector"],
                inplace=True,
                errors="ignore"
            )
            # Preserve column order and apply same styling as Admin
            all_cols = dfp_display.columns.tolist()
            styledp = dfp_display[all_cols].style \
                .applymap(style_status, subset=[c for c in ["Status"] if c in all_cols]) \
                .applymap(style_region, subset=[c for c in ["Region of operation"] if c in all_cols]) \
                .applymap(style_sector, subset=[c for c in ["Sector"] if c in all_cols])
            st.dataframe(styledp, use_container_width=True)

            # Detailed view and reject option for each portfolio project
            st.subheader("Project Details")
            for rec in filtered_portfolio:
                with st.expander(rec["Project registered name"]):
                    base_path = os.path.join(UPLOAD_FOLDER, rec["folder"])
                    summ, _ = load_layered_summary(base_path)
                    # Merge metadata
                    meta_display = {
                        "Project registered name": rec.get("Project registered name", rec.get("Project Name", "")),
                        "Date of incorporation": rec.get("Date of incorporation", ""),
                        "Primary sector / theme": rec.get("Primary sector / theme", ""),
                        "Headquarters country": rec.get("Headquarters country", ""),
                        "Contact e-mail": rec.get("Contact e-mail", ""),
                    }
                    combined = {**meta_display, **summ}
                    render_summary_grid(combined)
                    reject_key = f"reject_{role_tag}_{rec['folder']}"
                    if st.button("🚫 Reject", key=reject_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "eject")
                        st.warning("Project moved to Rejected.")
                        _rerun()

    # -------------------- Tab 3: Rejected -----------------------
    with tabs[2]:
        st.subheader(f"Rejected (_{len(rejected)} projects_)", divider="gray")
        if not rejected:
            st.info("No rejected projects.")
        else:
            for rec in rejected:
                with st.expander(rec["Project registered name"]):
                    base_path = os.path.join(UPLOAD_FOLDER, rec["folder"])
                    summ, _ = load_layered_summary(base_path)
                    render_summary_grid(summ)
                    reconsider_key = f"reconsider_{role_tag}_{rec['folder']}"
                    if st.button("↩️ Reconsider", key=reconsider_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "reconsider")
                        st.success("Moved back to Unclassified.")
                        _rerun()

elif is_portfolio_ncge or is_portfolio_ncgd:
    # ════════════════════════════════════════════════════════════
    # PORTFOLIO‑HOLDER DASHBOARD  (Step 3: UI skeleton)
    # ════════════════════════════════════════════════════════════
    role_tag = "NCGE" if is_portfolio_ncge else "NCGD"
    st.title(f"📂 Portfolio Holder – {role_tag}")

    # -------------- Gather basic project records --------------
    folders = [
        f for f in os.listdir(UPLOAD_FOLDER)
        if os.path.isdir(os.path.join(UPLOAD_FOLDER, f))
    ]
    records = []
    for fld in folders:
        base = os.path.join(UPLOAD_FOLDER, fld)
        info = {}
        meta_file = os.path.join(base, "info.txt")
        if os.path.exists(meta_file):
            for ln in open(meta_file).read().splitlines():
                if ":" in ln:
                    k, v = ln.split(":", 1)
                    info[k.strip()] = v.strip()
        # Layered summary
        summary_dict, _ = load_layered_summary(base)
        portfolio_raw = summary_dict.get("Portfolio", [])
        if isinstance(portfolio_raw, str):
            portfolio_list = [p.strip() for p in portfolio_raw.split(";") if p.strip()]
        elif isinstance(portfolio_raw, list):
            portfolio_list = portfolio_raw
        else:
            portfolio_list = []

        # Basic numeric parsing (for later analytics)
        try:
            rev_val = float(str(summary_dict.get("Last 12 months revenues (USD)", "0")).replace(",", ""))
        except:
            rev_val = 0.0
        try:
            irr_val = float(str(summary_dict.get("Expected IRR (%)", "0")).replace("%", ""))
        except:
            irr_val = 0.0

        records.append({
            "folder": fld,
            "Project": info.get("Project", fld),
            "Sector": info.get("Sector", ""),
            "Region": summary_dict.get("Region of operation", ""),
            "Status": summary_dict.get("Maturity stage", ""),
            "Portfolio": portfolio_list,
            "Revenues": rev_val,
            "IRR": irr_val,
        })

    # Split into unclassified vs in‑portfolio
    unclassified = [r for r in records if role_tag not in r["Portfolio"]]
    in_portfolio = [r for r in records if role_tag in r["Portfolio"]]

    # Two‑tab layout
    tab_uncls, tab_port = st.tabs(["🆕 Unclassified", "📂 My Portfolio"])

    # -------------------- Tab 1: Unclassified -------------------
    with tab_uncls:
        st.subheader("Unclassified projects")
        if not unclassified:
            st.success("All projects have been classified ✅")
        else:
            for rec in unclassified:
                with st.expander(rec["Project registered name"]):
                    # Reload full summary for detailed view
                    base_path = os.path.join(UPLOAD_FOLDER, rec["folder"])
                    summ, _ = load_layered_summary(base_path)

                    # Show full 3‑column grid (same helper as Admin)
                    render_summary_grid(summ)
                    st.markdown("---")

                    # Quick numeric highlights
                    col1, col2, col3 = st.columns(3)
                    col1.metric("Revenue (USD)", summ.get("Last 12 months revenues (USD)", "–"))
                    col2.metric("Financing Need (USD)", summ.get("Financing need or round size (USD)", "–"))
                    col3.metric("Expected IRR (%)", summ.get("Expected IRR (%)", "–"))

                    # Add/Eject buttons with persistence
                    add_key = f"add_{role_tag}_{rec['folder']}"
                    eject_key = f"eject_{role_tag}_{rec['folder']}"
                    colA, colB = st.columns(2)
                    if colA.button("➕ Add to portfolio", key=add_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "add")
                        st.success("Added to portfolio!")
                        _rerun()
                    if colB.button("❌ Eject", key=eject_key):
                        _set_portfolio_membership(rec["folder"], role_tag, "eject")
                        st.warning("Ejected.")
                        _rerun()
            st.caption("Click **Add** to move a project to your portfolio or **Eject** to remove it from view.")

    # -------------------- Tab 2: Portfolio view -----------------
    with tab_port:
        st.subheader(f"My Portfolio ({len(in_portfolio)} projects)")
        if not in_portfolio:
            st.info("No projects in your portfolio yet.")
        else:
            dfp = pd.DataFrame(in_portfolio)
            # Simple read‑only table; full analytics will be added later
            st.dataframe(
                dfp[["Project Name", "Sector", "Region", "Status", "Revenues", "IRR"]],
                use_container_width=True,
            )
            st.caption("Read‑only view. Editing and dashboards will arrive in later steps.")
else:
    # ═════════════════════════════════════════════════════════════
    # ENTREPRENEUR DASHBOARD  
    # ═════════════════════════════════════════════════════════════
    st.title("📥 Impact Project Room")

    # Mode selector
    mode = st.sidebar.radio("Mode", ["Submit New", "Edit Existing"])

    # ========== EDIT EXISTING SUBMISSION ========================
    if mode == "Edit Existing":
        st.subheader("Edit your existing submission")
        eid  = st.text_input("Project ID")
        epin = st.text_input("PIN", type="password")
        if st.button("Load Submission"):
            cred_file = os.path.join(UPLOAD_FOLDER, eid or "", "credentials.json")
            if os.path.exists(cred_file):
                creds = json.load(open(cred_file))
                if creds.get("pin") == epin:
                    # Load metadata, summary → session_state
                    base = os.path.join(UPLOAD_FOLDER, eid)
                    info = {}
                    for ln in open(os.path.join(base, "info.txt")).read().splitlines():
                        if ":" in ln:
                            k, v = ln.split(":", 1)
                            info[k.strip()] = v.strip()
                    summary = {}
                    try:
                        raw = open(os.path.join(base, "summary_entrepreneur.json")).read()
                        summary = json.loads(raw)["fields"] if "fields" in raw else json.loads(raw)
                    except Exception:
                        pass
                    st.session_state.form_meta = {
                        "Project": info.get("Project", ""),
                        "Incorporation date": info.get("Incorporation date", datetime.now().date()),
                        "Country HQ": info.get("Country HQ", ""),
                        "Sector": info.get("Sector", ""),
                        "Email": info.get("Email", ""),
                    }
                    st.session_state.form_files = []
                    st.session_state.form_summ = summary
                    st.session_state.edit_folder = eid
                    st.session_state.stage = "review"
                    _rerun()
                else:
                    st.error("Invalid PIN.")
            else:
                st.error("Invalid Project ID.")

    # Two‑stage form state machine
    if "stage" not in st.session_state:
        st.session_state.stage = "input"

    # ── Stage 1: initial inputs + AI generation ────────────────
    if mode == "Submit New" and st.session_state.stage == "input":

        with st.form("project_form"):
            proj     = st.text_input("Project registered name")
            inc_date = st.date_input("Date of incorporation")
            sector   = st.selectbox("Primary sector / theme", SECTOR_OPTIONS)
            country  = st.selectbox("Headquarters country", COUNTRY_OPTIONS)
            email    = st.text_input("Contact e‑mail")
            files    = st.file_uploader("Upload up to 5 files (any format)",
                                         accept_multiple_files=True)
            with st.expander("Read the full NDA terms"):
                st.markdown('''# NON-DISCLOSURE AGREEMENT
                            

By ticking the acceptance box, the undersigned party (“Accepting Party”) confirms agreement with the terms of this mutual non-disclosure agreement (“NDA”) entered into with KickImpact SA, Rue de Lyon 5, 1201 Genève – CHE-138.972.998 (“KISA”) (together referred to as the “Parties”, individually as a “Party”).

---

## 1. Confidential Information

“Confidential Information” includes:
- Any and all information disclosed by the DISCLOSING PARTY to the RECEIVING PARTY in any form (written, oral, electronic, visual), including without limitation: business plans, strategies, models, designs, processes, trademarks, inventions, financial data, forecasts, research, documentation, contacts, and any information related to the DISCLOSING PARTY’s activities.
- Any information obtained indirectly by the RECEIVING PARTY through analysis, review, or inspection of such data.
- Information from third parties held by the DISCLOSING PARTY and shared under this NDA.

“Intellectual Property Rights” include, but are not limited to, patents, copyrights, trademarks, trade secrets, design rights, software rights, and any other rights over confidential or proprietary material, whether registered or not.

---

## 2. Exclusions

Confidential Information shall not include information that:
- Was already public before disclosure;
- Becomes public through no fault of the RECEIVING PARTY;
- Can be evidenced as already possessed by the RECEIVING PARTY;
- Was lawfully obtained from a third party;
- Was independently developed by the RECEIVING PARTY;
- Is explicitly permitted for use by the DISCLOSING PARTY.

---

## 3. Ownership and Restrictions

- All Confidential Information remains the sole property of the DISCLOSING PARTY.
- No rights or licenses are granted by the disclosure of Confidential Information.

---

## 4–6. Obligations of the RECEIVING PARTY

The RECEIVING PARTY shall:
- Treat all Confidential Information with strict confidentiality;
- Not use it for any purpose other than the agreed Purpose;
- Not disclose it to third parties without prior written consent;
- Take all reasonable measures to prevent unauthorized access or use;
- Notify the DISCLOSING PARTY in case of any breach or suspected breach;
- Not attempt to circumvent or derive benefit from the Confidential Information;
- Not reproduce, store, or disseminate Confidential Information unless necessary to fulfill the Purpose.

---

## 7. Non-Solicitation

Neither Party shall solicit or attempt to solicit any employee, client, prospect, partner, or Board member of the other Party.

---

## 8. Legal Disclosures

If compelled by law to disclose Confidential Information, the concerned Party shall notify the other and limit disclosure to what is legally required.

---

## 9. No Obligation to Proceed

This NDA does not obligate either Party to pursue any specific transaction or relationship.

---

## 10. Public Statements

No press release or public communication regarding the relationship shall be made without prior written consent from the other Party.

---

## 11. Duration and Survival

- The NDA shall remain effective for two (2) years.
- Obligations of confidentiality survive indefinitely.
- Upon termination, all Confidential Information shall be returned or securely destroyed (with exceptions for backup and latent data).

---

## 12. No Warranty

Confidential Information is provided “as is”. The DISCLOSING PARTY makes no warranties and assumes no liability for the use or accuracy of the information.

---

## 13. Severability

If any provision of this NDA is deemed unenforceable, the remainder shall remain valid and enforceable to the fullest extent permitted by law.

---

## 14. Authority

Each Party warrants having the authority to enter into and comply with this NDA, which shall be binding and enforceable under its terms.

---

## 15. Assignment

Neither Party may assign this NDA without the prior written consent of the other. Any unauthorized assignment shall be void.

---

## 16. Governing Law and Jurisdiction

This NDA is governed by Swiss law. Any dispute shall be subject to the exclusive jurisdiction of the courts of Geneva, Switzerland, with right of appeal to the Swiss Supreme Court (Tribunal Fédéral).

---

**By ticking the box, you acknowledge that you have read, understood, and agreed to the terms of this Non-Disclosure Agreement.**
''')
            nda      = st.checkbox("I accept the NDA")
            st.info("By ticking the acceptance box, you confirm that you have read, understood, and agree to the terms of the Non-Disclosure Agreement.")

            generate_btn = st.form_submit_button("Generate AI Summary")

        if generate_btn:
            # Simple validation
            if not (proj and email and nda and files):
                st.warning("Fill required fields, upload at least one file, accept NDA.")
                st.stop()
            if len(files) > 5:
                st.warning("Max 5 files.")
                st.stop()

            # Extract text from first file & summarize
            from tempfile import NamedTemporaryFile
            first_file = files[0]
            suffix = Path(first_file.name).suffix
            with NamedTemporaryFile(delete=True, suffix=suffix) as tmp:
                tmp.write(first_file.read()); tmp.flush()
                text = extract_text_from_file(tmp.name)
            summary = summarize_project_with_gpt(text)

            # Stash in session_state & move to review stage
            st.session_state.form_meta  = {
                "Project": proj,
                "Incorporation date": inc_date,
                "Country HQ": country,
                "Sector": sector,
                "Email": email,
            }
            st.session_state.form_files = files
            st.session_state.form_summ  = summary
            st.session_state.stage = "review"
            _rerun()

    # ── Stage 2: review AI output, allow edits ─────────────────
    if st.session_state.stage == "review":
        st.header("Review & Edit Fields")
        edited = {}
        fs = st.session_state.form_summ  # shorthand
        # Edit initial entrepreneur inputs
        st.markdown("### Review Your Submission Details")
        # Primary sector / theme
        sector_val = st.session_state.form_meta.get("Sector", "")
        new_sector = st.selectbox(
            "Primary sector / theme",
            SECTOR_OPTIONS,
            index=SECTOR_OPTIONS.index(sector_val) if sector_val in SECTOR_OPTIONS else 0,
            key="review_sector"
        )
        st.session_state.form_meta["Sector"] = new_sector

        # Headquarters country
        country_val = st.session_state.form_meta.get("Country HQ", "")
        new_country = st.selectbox(
            "Headquarters country",
            COUNTRY_OPTIONS,
            index=COUNTRY_OPTIONS.index(country_val) if country_val in COUNTRY_OPTIONS else 0,
            key="review_country"
        )
        st.session_state.form_meta["Country HQ"] = new_country

        # Contact e-mail
        email_val = st.session_state.form_meta.get("Email", "")
        new_email = st.text_input(
            "Contact e-mail",
            value=email_val,
            key="review_email"
        )
        st.session_state.form_meta["Email"] = new_email

        st.markdown("---")

        st.markdown("### Review AI‑generated Summary")

        # Free‑text simple fields
        edited["Project Name"] = st.text_input(
            "Project Name", value=fs.get("Project Name", "")
        )
        edited["Business Model"] = st.text_input(
            "Business Model", value=fs.get("Business Model", "")
        )

        # Maturity stage dropdown
        ms_val = fs.get("Maturity stage", "")
        ms_idx = MATURITY_STAGES.index(ms_val) if ms_val in MATURITY_STAGES else 0
        edited["Maturity stage"] = st.selectbox(
            "Maturity stage", MATURITY_STAGES, index=ms_idx
        )

        # Region of operation – single select
        reg_options = list(REGION_COLORS.keys())
        reg_val = fs.get("Region of operation", "")
        reg_idx = reg_options.index(reg_val) if reg_val in reg_options else 0
        edited["Region of operation"] = st.selectbox(
            "Region of operation", reg_options, index=reg_idx
        )

        # Main country of current operations – single select
        mc_val = fs.get("Main country of current operations", "")
        mc_idx = COUNTRY_OPTIONS.index(mc_val) if mc_val in COUNTRY_OPTIONS else 0
        edited["Main country of current operations"] = st.selectbox(
            "Main country of current operations", COUNTRY_OPTIONS, index=mc_idx
        )

        # Instrument – single select
        inst_val = fs.get("Instrument", "")
        inst_idx = INSTRUMENT_OPTIONS.index(inst_val) if inst_val in INSTRUMENT_OPTIONS else 0
        edited["Instrument"] = st.selectbox(
            "Instrument", INSTRUMENT_OPTIONS, index=inst_idx
        )

        # Specific Sector(s) – multi-select
        ss_raw = fs.get("Specific Sector(s)", "")
        if isinstance(ss_raw, list):
            ss_defaults = [s for s in ss_raw if s in SPECIFIC_SECTOR_OPTIONS]
        else:
            ss_defaults = [s.strip() for s in str(ss_raw).split(";") if s.strip() in SPECIFIC_SECTOR_OPTIONS]
        edited["Specific Sector(s)"] = st.multiselect(
            "Specific Sector(s)", SPECIFIC_SECTOR_OPTIONS, default=ss_defaults
        )

        # Simple text inputs
        for fld in ["Core team", "Impact Area", "Key risks", "Barrier(s) to entry"]:
            edited[fld] = st.text_input(fld, value=fs.get(fld, ""))

        # Numeric helpers
        def _int_default(raw): return int(str(raw).replace(",", "")) if str(raw).replace(",","").isdigit() else 0
        def _float_default(raw):
            try: return float(str(raw).replace("%",""))
            except: return 0.0

        edited["Last 12 months revenues (USD)"] = st.number_input(
            "Last 12 months revenues (USD)", min_value=0,
            value=_int_default(fs.get("Last 12 months revenues (USD)", "")),
            step=1, format="%d"
        )
        edited["Market size or SOM (USD)"] = st.number_input(
            "Market size or SOM (USD)", min_value=0,
            value=_int_default(fs.get("Market size or SOM (USD)", "")),
            step=1, format="%d"
        )
        edited["Financing need or round size (USD)"] = st.number_input(
            "Financing need or round size (USD)", min_value=0,
            value=_int_default(fs.get("Financing need or round size (USD)", "")),
            step=1, format="%d"
        )
        
        edited["Breakeven year"] = st.number_input(
            "Breakeven year",
            min_value=1900,
            max_value=2100,
            value=safe_int(fs.get("Breakeven year"), datetime.now().year),
            step=1,
        )

        edited["Expected IRR (%)"] = st.number_input(
            "Expected IRR (%)", min_value=0.0, max_value=100.0,
            value=_float_default(fs.get("Expected IRR (%)", "")),
            step=0.1, format="%.2f"
        )
        edited["Use of proceeds"] = st.text_input(
            "Use of proceeds",
            value=fs.get("Use of proceeds", ""),
            help="e.g. ‘90% AI investment, 5% salaries, 5% admin’ / 'Funds will be needed to buy land and proceed to Bamboo planting. Funds will be used to cover OPEX and CAPEX, until next tranche of capital become available for carbon credit clients in 2025."
        )

        # SDGs multiselect (max 3)
        raw_sdgs = fs.get("Main SDGs targeted (3 max)", "")
        raw_list = [s.strip() for s in raw_sdgs.split(";")] if isinstance(raw_sdgs, str) else raw_sdgs
        default_sdgs = [s for s in _match_sdgs(raw_list) if s in SDG_OPTIONS]
        if not default_sdgs:
            default_sdgs = []
        edited["Main SDGs targeted (3 max)"] = st.multiselect(
            "Main SDGs targeted (3 max)", SDG_OPTIONS, default=default_sdgs, max_selections=3
        )

        # Long‑text fields
        for fld in ["Problem", "Solution"]:
            edited[fld] = st.text_area(fld, value=fs.get(fld, ""))

        # ---------- Submit button ----------
        if st.button("✅ Confirm & Submit"):
            if "edit_folder" in st.session_state:
                _update_submission(
                    st.session_state.edit_folder,
                    st.session_state.form_meta,
                    edited,
                )
                fld = os.path.join(UPLOAD_FOLDER, st.session_state.edit_folder)
                token = None
            else:
                fld, token = _save_submission(
                    st.session_state.form_meta,
                    st.session_state.form_files,
                    edited,
                )
            # Email alert
            email_admin(
                "📥 New Impact Project Room submission",
                f"Project: {edited.get('Project Name', st.session_state.form_meta['Project'])}\n"
                f"Sector:  {st.session_state.form_meta['Sector']}\n"
                f"Country: {st.session_state.form_meta['Country HQ']}\n"
                f"Uploaded: {datetime.now().isoformat(timespec='seconds')}"
            )
            # Stash details & move on
            st.session_state.submitted_details = edited
            st.session_state._last_submission_fld = fld
            st.session_state._last_submission_token = token
            st.session_state.stage = "done"
            _rerun()

    # ── Stage 3: confirmation screen ─────────────────────────────
    if st.session_state.stage == "done":
        st.title("🎉 Submission Complete")
        st.write("Thank you for the update! You can now save your Project ID and PIN to edit later.")
        st.subheader("Your submitted details:")
        # Show both original inputs and final AI-edited summary
        combined = {**st.session_state.form_meta, **st.session_state.submitted_details}
        for k, v in combined.items():
            st.markdown(f"- **{k}:** {v}")
        if st.session_state._last_submission_token:
            st.subheader("Your Project ID and PIN:")
            proj_id = os.path.basename(st.session_state._last_submission_fld)
            st.markdown(f"- **Project ID:** `{proj_id}`")
            st.markdown(f"- **PIN:** `{st.session_state._last_submission_token}`")
            st.info("Please save your Project ID and PIN to edit your submission later.")
        if st.button("Submit another project"):
            for k in (
                "stage","form_meta","form_files","form_summ","submitted_details",
                "_last_submission_fld","_last_submission_token","edit_folder"
            ):
                st.session_state.pop(k, None)
            _rerun()

# ===================== EXPORT BLOCK (Admin) =======================
if is_admin:
    # ... (other admin code above)
    # EXPORT to CSV/Excel (add Portfolio field)
    if st.button("Export projects to CSV/Excel", key="export_projects"):
        rows = []
        for rec in records:
            # Build row from AI_FIELDS
            parsed = load_layered_summary(os.path.join(UPLOAD_FOLDER, rec["folder"]))[0]
            row = {f: "" for f in AI_FIELDS}
            row["Portfolio"] = "; ".join(parsed.get("Portfolio", [])) if isinstance(parsed.get("Portfolio", []), list) else parsed.get("Portfolio", "")
            row["Project"] = rec["Project"]
            row["Country"] = rec["Country HQ"]
            row["Sector"] = rec["Sector"]
            row["Status"] = rec["Status"]
            row["Email"] = rec["Email"]
            for f in AI_FIELDS:
                row[f] = parsed.get(f, "")
            rows.append(row)
        df = pd.DataFrame(rows).reindex(columns=["Project", "Country", "Sector", "Status", "Portfolio", "Email"] + AI_FIELDS)
        # Export as CSV/Excel
        csv = df.to_csv(index=False)
        st.download_button("Download CSV", csv, file_name="impact_projects.csv", mime="text/csv")
        excel = df.to_excel(index=False, engine="openpyxl")
        st.download_button("Download Excel", excel, file_name="impact_projects.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
